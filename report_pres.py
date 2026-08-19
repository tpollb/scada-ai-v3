# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Цветовая схема — тёмная тема
BG_DARK = RGBColor(18, 18, 18)           # #121212 — основной фон
BG_PANEL = RGBColor(28, 28, 28)          # #1C1C1C — панели/карточки
BG_LIGHTER = RGBColor(38, 38, 38)        # #262626 — акцентные блоки
TEXT_PRIMARY = RGBColor(240, 240, 240)   # #F0F0F0 — основной текст
TEXT_SECONDARY = RGBColor(170, 170, 170) # #AAAAAA — вторичный текст
TEXT_DIM = RGBColor(120, 120, 120)       # #787878 — приглушённый текст
ACCENT_BLUE = RGBColor(0, 120, 212)      # #0078D4 — акцент синий
ACCENT_LIGHT = RGBColor(79, 165, 230)    # #4FA5E6 — светлый акцент
ACCENT_GREEN = RGBColor(76, 175, 80)     # #4CAF50 — позитив/экономика
ACCENT_AMBER = RGBColor(255, 179, 0)     # #FFB300 — предупреждения/важно
LINE_COLOR = RGBColor(55, 55, 55)        # #373737 — разделители

FONT_NAME = 'Segoe UI Light'
FONT_BOLD = 'Segoe UI Semibold'

def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text,
                font_size=14, color=TEXT_PRIMARY, bold=False,
                alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_panel(slide, left, top, width, height, color=BG_PANEL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_line(slide, left, top, width, color=LINE_COLOR, thickness=1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, height, color=ACCENT_BLUE, width=0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def bullet_list(tf, items, font_size=13, color=TEXT_PRIMARY, bullet_color=ACCENT_LIGHT):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"›  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)

# =========================================================
# СЛАЙД 1: ТИТУЛЬНЫЙ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Верхняя тонкая линия-акцент
add_line(slide, 1.0, 1.5, 3.0, ACCENT_BLUE, 3)

# Заголовок
add_textbox(slide, 1.0, 1.8, 11.0, 1.2,
            "СИСТЕМА ОТЧЁТОВ",
            font_size=48, color=TEXT_PRIMARY, bold=False, font_name=FONT_NAME)

# Подзаголовок
add_textbox(slide, 1.0, 2.9, 11.0, 0.8,
            "iRidi SCADA-BMS  ·  Stimulsoft Reports Engine",
            font_size=22, color=ACCENT_LIGHT, bold=False, font_name=FONT_NAME)

# Разделитель
add_line(slide, 1.0, 3.9, 11.0, LINE_COLOR, 1)

# Описание
add_textbox(slide, 1.0, 4.2, 9.0, 1.5,
            "Формирование, визуализация и экспорт отчётности "
            "по инженерным системам зданий и сооружений. "
            "Инструмент контроля, аудита и экономической оптимизации.",
            font_size=16, color=TEXT_SECONDARY, bold=False, font_name=FONT_NAME)

# Нижний блок с метаинформацией
add_panel(slide, 1.0, 6.0, 11.3, 0.9, BG_PANEL)
add_textbox(slide, 1.4, 6.15, 10.5, 0.6,
            "Платформы: Windows · Linux · Debian 12    |    "
            "Экспорт: PDF · Excel · Word · HTML · CSV · JSON · XML    |    "
            "Источники: PostgreSQL · OData · REST API",
            font_size=12, color=TEXT_DIM, bold=False, font_name=FONT_NAME)

# =========================================================
# СЛАЙД 2: ОБЩЕЕ ОПИСАНИЕ СИСТЕМЫ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ОБЩЕЕ ОПИСАНИЕ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Архитектура модуля отчётности",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_BLUE, 3)

# Левая панель — суть системы
add_panel(slide, 1.0, 2.1, 5.8, 4.6, BG_PANEL)
add_accent_bar(slide, 1.0, 2.1, 4.6, ACCENT_BLUE)
add_textbox(slide, 1.3, 2.3, 5.3, 0.4,
            "ЧТО ЭТО", font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 1.3, 2.65, 5.3, 1.2,
            "Встроенный модуль генерации отчётов на базе движка "
            "Stimulsoft Reports.JS, интегрированный в iRidi SCADA-BMS. "
            "Работает на одном сервере со SCADA Server, подключается "
            "к базе данных PostgreSQL и формирует отчёты произвольной сложности.",
            font_size=14, color=TEXT_PRIMARY)

add_textbox(slide, 1.3, 4.0, 5.3, 0.4,
            "КЛЮЧЕВЫЕ ОСОБЕННОСТИ", font_size=11, color=ACCENT_LIGHT)
tf = add_textbox(slide, 1.3, 4.35, 5.3, 2.2, "", font_size=13)
bullet_list(tf, [
    "Серверная генерация — отчёты строятся на сервере, не на клиенте",
    "Веб-интерфейс — редактирование шаблонов прямо из браузера",
    "Шаблоны — готовые формы «История событий» и «Значения тегов»",
    "Группировка — по тегам, датам (день / неделя / месяц), зонам",
    "Агрегация — суммы, средние, итоги по группам через Totals",
    "Фильтрация — по произвольным полям, в т.ч. не включённым в отчёт",
    "Экспорт — PDF, Excel, Word, HTML, CSV, JSON, SVG, XML и др."
])

# Правая панель — принципы работы
add_panel(slide, 7.1, 2.1, 5.2, 2.1, BG_PANEL)
add_accent_bar(slide, 7.1, 2.1, 2.1, ACCENT_GREEN)
add_textbox(slide, 7.4, 2.3, 4.7, 0.4,
            "ПРИНЦИП РАБОТЫ", font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 7.4, 2.65, 4.7, 1.5,
            "SCADA Server собирает данные с контроллеров и сохраняет "
            "в PostgreSQL. Модуль Stimulsoft подключается к этой же БД "
            "и по запросу формирует отчёт по шаблону. Результат — "
            "готовый документ в выбранном формате.",
            font_size=13, color=TEXT_PRIMARY)

add_panel(slide, 7.1, 4.5, 5.2, 2.2, BG_PANEL)
add_accent_bar(slide, 7.1, 4.5, 2.2, ACCENT_AMBER)
add_textbox(slide, 7.4, 4.7, 4.7, 0.4,
            "ТИПЫ СОБЫТИЙ В ОТЧЁТАХ", font_size=11, color=ACCENT_LIGHT)
tf = add_textbox(slide, 7.4, 5.05, 4.7, 1.6, "", font_size=12)
bullet_list(tf, [
    "UserLogin / UserLogout — вход и выход пользователей",
    "Alarm / AlarmConfirmed / AlarmsCleared — аварийные события",
    "TagAccessed / TagChanged — изменения значений тегов",
    "CustomEvent / CustomAlarm — пользовательские события",
    "DbConnected / DbDisconnected — состояние БД",
    "LogsCleaningPerformed — очистка журналов"
], font_size=11)

# Нижний слайд-подвал
add_textbox(slide, 1.0, 6.9, 11.0, 0.4,
            "Модуль входит в стандартную поставку iRidi SCADA-BMS. "
            "Лицензия Stimulsoft — royalty-free, без дополнительных отчислений.",
            font_size=11, color=TEXT_DIM)

# =========================================================
# СЛАЙД 3: ЧТО ДАЁТ ПОЛЬЗОВАТЕЛЮ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ЦЕННОСТЬ ДЛЯ ПОЛЬЗОВАТЕЛЯ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Что получает заказчик системы",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_BLUE, 3)

# 4 карточки преимуществ
cards = [
    ("ПРОЗРАЧНОСТЬ",
     "Полная картина работы инженерных "
     "систем: от состояния датчиков до "
     "действий операторов. Каждое событие "
     "зафиксировано с точностью до секунды.",
     ACCENT_BLUE),
    ("ДОКАЗАТЕЛЬНОСТЬ",
     "Юридически значимые документы: "
     "отчёты с timestamp, экспорт в PDF "
     "с подписями. Основа для споров "
     "с ресурсоснабжающими организациями.",
     ACCENT_GREEN),
    ("АВТОМАТИЗАЦИЯ",
     "Шаблоны отчётов формируются "
     "автоматически. Ежедневные, еженедельные "
     "и ежемесячные сводки — без ручного "
     "вмешательства оператора.",
     ACCENT_AMBER),
    ("ИНТЕГРАЦИЯ",
     "Готовые форматы для бухгалтерии, "
     "руководства, контролирующих органов. "
     "Экспорт в Excel для анализа, "
     "в PDF для архива и отчётности.",
     ACCENT_LIGHT),
]

for i, (title, desc, color) in enumerate(cards):
    x = 1.0 + i * 2.95
    add_panel(slide, x, 2.1, 2.75, 3.8, BG_PANEL)
    add_accent_bar(slide, x, 2.1, 3.8, color)
    add_textbox(slide, x + 0.25, 2.35, 2.3, 0.4,
                title, font_size=13, color=color)
    add_textbox(slide, x + 0.25, 2.75, 2.3, 3.0,
                desc, font_size=13, color=TEXT_PRIMARY)

# Нижний блок — итоговое позиционирование
add_panel(slide, 1.0, 6.2, 11.3, 0.9, BG_PANEL)
add_accent_bar(slide, 1.0, 6.2, 0.9, ACCENT_BLUE)
add_textbox(slide, 1.3, 6.35, 10.7, 0.6,
            "Итог: система отчётов трансформирует сырые данные SCADA "
            "в управленческие, финансовые и юридические документы — "
            "без привлечения программистов и внешних BI-систем.",
            font_size=14, color=TEXT_PRIMARY)

# =========================================================
# СЛАЙД 4: СИТУАЦИИ ПРИМЕНЕНИЯ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "СИТУАЦИИ ПРИМЕНЕНИЯ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Когда и как использовать отчёты",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_BLUE, 3)

scenarios = [
    ("ЭКСПЛУАТАЦИОННЫЙ КОНТРОЛЬ",
     "Ежесменный отчёт оператора: состояние систем, сработавшие "
     "аварии, подтверждённые события, действия персонала. "
     "Группировка по зонам и времени суток.",
     "Главный инженер, диспетчер, начальник смены"),
    ("АВАРИЙНЫЙ АНАЛИЗ",
     "Отчёт по инцидентам за период: последовательность событий, "
     "время реакции, длительность аварий. Выявление системных "
     "проблем и повторяющихся сбоев.",
     "Служба эксплуатации, сервисная служба"),
    "ЭНЕРГЕТИЧЕСКИЙ АУДИТ",
    "Значения тегов счётчиков и датчиков с группировкой "
    "по дням / неделям / месяцам. Суммарное потребление, "
    "пиковые значения, динамика по каждой зоне объекта.",
    "Энергоменеджер, главный энергетик, бухгалтерия",
    "СВЕРКА С РЕСУРСОСНАБЖАЮЩИМИ ОРГАНИЗАЦИЯМИ",
    "Формирование отчётов по показаниям приборов учёта "
    "за расчётный период. Экспорт в Excel и PDF для передачи "
    "в энергосбыт, водоканал, теплосеть.",
    "Бухгалтерия, финансовый отдел, юристы",
    "АУДИТ ДЕЙСТВИЙ ПЕРСОНАЛА",
    "Отчёт по входам, выходам и действиям пользователей. "
    "Фиксация несанкционированного доступа, изменение "
    "уставок, подтверждение аварий.",
    "Служба безопасности, руководство объекта",
    "РЕГЛАМЕНТНАЯ ОТЧЁТНОСТЬ",
    "Автоматические ежемесячные сводки для руководства, "
    "собственников, арендаторов. KPI эксплуатации, "
    "статистика отказов, потребление ресурсов.",
    "Управляющая компания, собственник, арендодатель"
]

# Переупаковка: 3 карточки в ряд, 2 ряда
cards_data = [
    ("ЭКСПЛУАТАЦИОННЫЙ КОНТРОЛЬ",
     "Ежесменный отчёт: состояние систем, аварии, "
     "действия персонала. Группировка по зонам и "
     "времени суток.",
     "Диспетчер · Главный инженер",
     ACCENT_BLUE),
    ("АВАРИЙНЫЙ АНАЛИЗ",
     "Последовательность инцидентов, время реакции, "
     "длительность аварий. Выявление системных "
     "проблем и повторяющихся сбоев.",
     "Служба эксплуатации · Сервис",
     ACCENT_AMBER),
    ("ЭНЕРГЕТИЧЕСКИЙ АУДИТ",
     "Показания счётчиков и датчиков по дням, "
     "неделям, месяцам. Пики, динамика, "
     "суммарное потребление по зонам.",
     "Энергоменеджер · Бухгалтерия",
     ACCENT_GREEN),
    ("СВЕРКА С РСО",
     "Показания приборов учёта за период. "
     "Экспорт в Excel и PDF для передачи "
     "в энергосбыт, водоканал, теплосеть.",
     "Бухгалтерия · Финансовый отдел",
     ACCENT_LIGHT),
    ("АУДИТ ПЕРСОНАЛА",
     "Входы, выходы, действия операторов. "
     "Фиксация несанкционированного доступа "
     "и изменения уставок.",
     "СБ · Руководство объекта",
     ACCENT_BLUE),
    ("РЕГЛАМЕНТНАЯ ОТЧЁТНОСТЬ",
     "Ежемесячные сводки для руководства "
     "и собственников. KPI, статистика "
     "отказов, потребление ресурсов.",
     "УК · Собственник · Арендодатель",
     ACCENT_AMBER),
]

for i, (title, desc, who, color) in enumerate(cards_data):
    row = i // 3
    col = i % 3
    x = 1.0 + col * 3.85
    y = 2.1 + row * 2.55

    add_panel(slide, x, y, 3.65, 2.35, BG_PANEL)
    add_accent_bar(slide, x, y, 2.35, color)
    add_textbox(slide, x + 0.25, y + 0.2, 3.2, 0.35,
                title, font_size=12, color=color)
    add_textbox(slide, x + 0.25, y + 0.55, 3.2, 1.2,
                desc, font_size=12, color=TEXT_PRIMARY)
    add_textbox(slide, x + 0.25, y + 1.85, 3.2, 0.35,
                who, font_size=10, color=TEXT_DIM)

# =========================================================
# СЛАЙД 5: ЭКОНОМИЧЕСКИЙ ЭФФЕКТ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ЭКОНОМИЧЕСКИЙ ЭФФЕКТ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Измеримая выгода для собственника и УК",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_GREEN, 3)

# 4 карточки экономического эффекта
eco_cards = [
    ("15–25%",
     "снижение затрат на энергоресурсы за счёт "
     "выявления аномалий потребления и неэффективных "
     "режимов работы оборудования",
     "Энергоаудит через отчёты «Значения тегов» "
     "с группировкой по периодам"),
    ("30–50%",
     "сокращение трудозатрат на подготовку "
     "регламентной отчётности для руководства, "
     "собственников и контролирующих органов",
     "Автоматические шаблоны вместо ручного "
     "сбора данных из журналов"),
    ("100%",
     "возврат переплат при успешном оспаривании "
     "начислений ресурсоснабжающих организаций "
     "на основании данных SCADA",
     "Отчёты с timestamp и подписью как "
     "доказательная база в спорах с РСО"),
    ("0 ₽",
     "дополнительных отчислений за формирование "
     "отчётов конечными пользователями — "
     "лицензия Stimulsoft royalty-free",
     "Модуль входит в поставку iRidi SCADA-BMS, "
     "не требует отдельных лицензий"),
]

for i, (num, desc, how) in enumerate(eco_cards):
    x = 1.0 + (i % 2) * 5.8
    y = 2.1 + (i // 2) * 2.4

    add_panel(slide, x, y, 5.6, 2.2, BG_PANEL)
    # Крупная цифра слева
    add_textbox(slide, x + 0.2, y + 0.25, 1.5, 0.8,
                num, font_size=34, color=ACCENT_GREEN,
                font_name=FONT_NAME)
    # Описание справа
    add_textbox(slide, x + 1.8, y + 0.25, 3.5, 1.0,
                desc, font_size=13, color=TEXT_PRIMARY)
    # Как достигается
    add_textbox(slide, x + 1.8, y + 1.4, 3.5, 0.7,
                how, font_size=11, color=TEXT_DIM)

# Итоговая панель
add_panel(slide, 1.0, 6.9, 11.3, 0.5, BG_PANEL)
add_textbox(slide, 1.3, 6.95, 10.7, 0.4,
            "Типовая окупаемость модуля отчётности: 3–6 месяцев "
            "только за счёт оспаривания начислений РСО и оптимизации энергопотребления.",
            font_size=12, color=TEXT_SECONDARY)

# =========================================================
# СЛАЙД 6: ОСПАРИВАНИЕ ДАННЫХ ЭНЕРГОСБЫТОВЫХ КОМПАНИЙ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ОСПАРИВАНИЕ НАЧИСЛЕНИЙ РСО",
            font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Отчёт SCADA как доказательная база",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_AMBER, 3)

# Левая панель — проблема
add_panel(slide, 1.0, 2.1, 5.6, 4.6, BG_PANEL)
add_accent_bar(slide, 1.0, 2.1, 4.6, ACCENT_AMBER)
add_textbox(slide, 1.3, 2.3, 5.1, 0.4,
            "ТИПОВЫЕ СИТУАЦИИ СПОРОВ С РСО",
            font_size=12, color=ACCENT_AMBER)

disputes = [
    "Расхождение показаний: счётчик РСО и прибор учёта "
    "на стороне потребителя дают разные значения",
    "Начисления по нормативу при фактическом наличии "
    "работающего прибора учёта",
    "Завышение пикового потребления из-за ошибок "
    "интеграции или неверных коэффициентов",
    "Начисления за периоды, когда объект фактически "
    "не потреблял ресурс (простои, аварии на стороне РСО)",
    "Некорректная тарификация: применение повышенного "
    "тарифа к объёмам, попадающим под базовый"
]
tf = add_textbox(slide, 1.3, 2.75, 5.1, 3.8, "", font_size=13)
bullet_list(tf, disputes, font_size=12)

# Правая панель — как SCADA решает
add_panel(slide, 6.9, 2.1, 5.4, 4.6, BG_PANEL)
add_accent_bar(slide, 6.9, 2.1, 4.6, ACCENT_GREEN)
add_textbox(slide, 7.2, 2.3, 4.9, 0.4,
            "ЧТО ПРЕДОСТАВЛЯЕТ SCADA-BMS",
            font_size=12, color=ACCENT_GREEN)

add_textbox(slide, 7.2, 2.75, 4.9, 0.4,
            "1. НЕПРЕРЫВНАЯ ФИКСАЦИЯ", font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 7.2, 3.05, 4.9, 0.5,
            "Архив значений тегов с точностью до секунды. "
            "Невозможно оспорить: данные хранятся в PostgreSQL, "
            "изменения фиксируются в журнале.",
            font_size=12, color=TEXT_PRIMARY)

add_textbox(slide, 7.2, 3.65, 4.9, 0.4,
            "2. ГРУППИРОВКА ПО ПЕРИОДАМ", font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 7.2, 3.95, 4.9, 0.5,
            "Отчёт за конкретный расчётный период РСО. "
            "Функция DayOfYear() и группировка по тегам "
            "формируют точную картину потребления.",
            font_size=12, color=TEXT_PRIMARY)

add_textbox(slide, 7.2, 4.55, 4.9, 0.4,
            "3. ЮРИДИЧЕСКИ ЗНАЧИМЫЙ ФОРМАТ",
            font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 7.2, 4.85, 4.9, 0.5,
            "Экспорт в PDF с метаданными, подписями, "
            "временными метками. Excel — для сверки "
            "построчно с актами РСО.",
            font_size=12, color=TEXT_PRIMARY)

add_textbox(slide, 7.2, 5.45, 4.9, 0.4,
            "4. ИСТОРИЯ СОБЫТИЙ", font_size=11, color=ACCENT_LIGHT)
add_textbox(slide, 7.2, 5.75, 4.9, 0.5,
            "Логи аварий, отключений, простоев. "
            "Подтверждение периодов, когда потребление "
            "объективно отсутствовало.",
            font_size=12, color=TEXT_PRIMARY)

# Нижняя панель — процесс
add_panel(slide, 1.0, 6.85, 11.3, 0.55, BG_PANEL)
add_textbox(slide, 1.3, 6.93, 10.7, 0.4,
            "Процесс:  запрос РСО  →  формирование отчёта в SCADA  →  "
            "экспорт PDF/Excel  →  приложение к претензии / исковому заявлению",
            font_size=12, color=TEXT_SECONDARY)

# =========================================================
# СЛАЙД 7: ИНТЕГРАЦИЯ С БУХГАЛТЕРИЕЙ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ИНТЕГРАЦИЯ С БУХГАЛТЕРИЕЙ",
            font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Сверка с ресурсоснабжающими организациями",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_BLUE, 3)

# Этапы процесса (горизонтальные карточки)
steps = [
    ("01", "СБОР ДАННЫХ",
     "SCADA Server автоматически опрашивает "
     "приборы учёта и сохраняет значения "
     "тегов в PostgreSQL с привязкой к "
     "дате, времени и зоне объекта."),
    ("02", "ФОРМИРОВАНИЕ ОТЧЁТА",
     "Бухгалтер или энергоменеджер "
     "запускает шаблон «Значения тегов» "
     "с группировкой по месяцам и "
     "фильтром по нужным приборам учёта."),
    ("03", "АГРЕГАЦИЯ И ЭКСПОРТ",
     "Функция Sum() и итоги по группам "
     "формируют суммарные значения. "
     "Экспорт в Excel для сверки с "
     "актами РСО, в PDF — для архива."),
    ("04", "СВЕРКА С РСО",
     "Построчное сравнение данных SCADA "
     "с актами ресурсоснабжающих компаний. "
     "Выявление расхождений, формирование "
     "претензии с приложением отчёта."),
]

for i, (num, title, desc) in enumerate(steps):
    x = 1.0 + i * 3.05
    add_panel(slide, x, 2.1, 2.85, 3.1, BG_PANEL)
    add_textbox(slide, x + 0.2, 2.25, 0.6, 0.5,
                num, font_size=28, color=ACCENT_BLUE,
                font_name=FONT_NAME)
    add_textbox(slide, x + 0.9, 2.35, 1.8, 0.4,
                title, font_size=13, color=ACCENT_LIGHT)
    add_textbox(slide, x + 0.2, 2.9, 2.45, 2.1,
                desc, font_size=12, color=TEXT_PRIMARY)

# Блок: форматы передачи в бухгалтерию
add_panel(slide, 1.0, 5.45, 11.3, 1.7, BG_PANEL)
add_accent_bar(slide, 1.0, 5.45, 1.7, ACCENT_BLUE)
add_textbox(slide, 1.3, 5.6, 10.7, 0.4,
            "ФОРМАТЫ ПЕРЕДАЧИ В БУХГАЛТЕРИЮ",
            font_size=12, color=ACCENT_LIGHT)

formats = [
    ("Excel (.xls)",
     "Основной формат для сверки. "
     "Структурированные таблицы с "
     "группировкой и итогами. "
     "Совместим с 1С и SAP."),
    ("PDF (.pdf)",
     "Формат для архива и официальных "
     "документов. Не редактируется, "
     "поддерживает подписи и шифрование. "
     "Юридическая значимость."),
    ("CSV / JSON / XML",
     "Машиночитаемые форматы для "
     "автоматической загрузки в "
     "учётные системы и ERP. "
     "Интеграция через REST API."),
    ("HTML",
     "Интерактивные отчёты для "
     "внутреннего портала. "
     "Просмотр в браузере без "
     "дополнительного ПО."),
]

for i, (fmt, desc) in enumerate(formats):
    x = 1.3 + i * 2.75
    add_textbox(slide, x, 6.0, 2.5, 0.3,
                fmt, font_size=13, color=ACCENT_GREEN)
    add_textbox(slide, x, 6.3, 2.5, 0.8,
                desc, font_size=11, color=TEXT_SECONDARY)

# =========================================================
# СЛАЙД 8: ТЕХНИЧЕСКИЕ ВОЗМОЖНОСТИ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 0.6, 11.0, 0.6,
            "ТЕХНИЧЕСКИЕ ВОЗМОЖНОСТИ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 0.95, 11.0, 0.7,
            "Платформы, источники данных, форматы экспорта",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 1.7, 2.5, ACCENT_BLUE, 3)

# Три колонки
# 1. Платформы
add_panel(slide, 1.0, 2.1, 3.6, 4.6, BG_PANEL)
add_accent_bar(slide, 1.0, 2.1, 4.6, ACCENT_BLUE)
add_textbox(slide, 1.25, 2.3, 3.1, 0.4,
            "ПЛАТФОРМЫ РАЗВЁРТЫВАНИЯ",
            font_size=11, color=ACCENT_LIGHT)
tf = add_textbox(slide, 1.25, 2.7, 3.1, 3.8, "", font_size=12)
bullet_list(tf, [
    "Windows (x86/x64)",
    "Linux (Debian 12, iRidi Server)",
    "iRidi SCADA-BMS Сервер кластера",
    "Node.js — серверная генерация",
    "Браузер — веб-интерфейс дизайнера",
    "Electron — настольный дизайнер",
    "Сервис устанавливается на одном "
    "устройстве со SCADA Server",
    "PostgreSQL — может быть удалённой"
], font_size=11)

# 2. Источники данных
add_panel(slide, 4.85, 2.1, 3.6, 4.6, BG_PANEL)
add_accent_bar(slide, 4.85, 2.1, 4.6, ACCENT_GREEN)
add_textbox(slide, 5.1, 2.3, 3.1, 0.4,
            "ИСТОЧНИКИ ДАННЫХ",
            font_size=11, color=ACCENT_LIGHT)
tf = add_textbox(slide, 5.1, 2.7, 3.1, 3.8, "", font_size=12)
bullet_list(tf, [
    "PostgreSQL — основная БД SCADA",
    "MS SQL, MySQL, Oracle, Firebird",
    "Mongo DB, SQLite",
    "OData (все стандартные типы)",
    "REST API — внешние сервисы",
    "JSON, XML, CSV, Excel файлы",
    "ODBC / PDO — любые СУБД",
    "Теги SCADA: значения и события"
], font_size=11)

# 3. Форматы экспорта
add_panel(slide, 8.7, 2.1, 3.6, 4.6, BG_PANEL)
add_accent_bar(slide, 8.7, 2.1, 4.6, ACCENT_AMBER)
add_textbox(slide, 8.95, 2.3, 3.1, 0.4,
            "ФОРМАТЫ ЭКСПОРТА",
            font_size=11, color=ACCENT_LIGHT)
tf = add_textbox(slide, 8.95, 2.7, 3.1, 3.8, "", font_size=12)
bullet_list(tf, [
    "PDF — с поддержкой подписей, PDF/A",
    "Excel (.xls) — для 1С и сверок",
    "Word (.doc) — текстовые документы",
    "HTML — веб-представление",
    "CSV — для импорта в ERP",
    "JSON / XML — для интеграций",
    "SVG / PNG / TIFF — графика",
    "RTF, DBF, ODS и другие"
], font_size=11)

# Нижний блок — ИИ-функции
add_panel(slide, 1.0, 6.85, 11.3, 0.55, BG_PANEL)
add_textbox(slide, 1.3, 6.93, 10.7, 0.4,
            "ИИ-ассистент в дизайнере: автоматическая проверка соединений к БД, "
            "анализ выражений, конвертация сценариев Blockly, исправление кода.",
            font_size=12, color=TEXT_SECONDARY)

# =========================================================
# СЛАЙД 9: ИТОГ
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.0, 1.2, 11.0, 0.6,
            "ИТОГИ", font_size=12, color=ACCENT_LIGHT)
add_textbox(slide, 1.0, 1.55, 11.0, 0.7,
            "Система отчётов iRidi SCADA-BMS",
            font_size=32, color=TEXT_PRIMARY)
add_line(slide, 1.0, 2.3, 2.5, ACCENT_BLUE, 3)

# 5 ключевых тезисов
theses = [
    ("01",
     "Прозрачность эксплуатации",
     "Полная фиксация всех событий инженерных систем "
     "с точностью до секунды. Данные невозможно подделать "
     "или оспорить — они хранятся в PostgreSQL."),
    ("02",
     "Экономическая эффективность",
     "Снижение затрат на энергоресурсы на 15–25%, "
     "сокращение трудозатрат на отчётность на 30–50%, "
     "возврат переплат при спорах с РСО."),
    ("03",
     "Юридическая защита",
     "Отчёты в PDF с метаданными — доказательная база "
     "в спорах с ресурсоснабжающими организациями, "
     "приложения к претензиям и исковым заявлениям."),
    ("04",
     "Интеграция без программистов",
     "Веб-интерфейс дизайнера, готовые шаблоны, "
     "экспорт во все популярные форматы. Бухгалтерия "
     "получает документы без привлечения IT-отдела."),
    ("05",
     "Отсутствие дополнительных затрат",
     "Модуль входит в стандартную поставку iRidi SCADA-BMS. "
     "Лицензия Stimulsoft — royalty-free, без отчислений "
     "за каждого пользователя или отчёт."),
]

for i, (num, title, desc) in enumerate(theses):
    y = 2.6 + i * 0.95
    add_panel(slide, 1.0, y, 11.3, 0.85, BG_PANEL)
    add_textbox(slide, 1.2, y + 0.15, 0.6, 0.55,
                num, font_size=22, color=ACCENT_BLUE,
                font_name=FONT_NAME)
    add_textbox(slide, 1.9, y + 0.15, 2.8, 0.55,
                title, font_size=14, color=TEXT_PRIMARY)
    add_textbox(slide, 4.7, y + 0.15, 7.4, 0.6,
                desc, font_size=12, color=TEXT_SECONDARY)

# Подвал
add_textbox(slide, 1.0, 7.0, 11.3, 0.4,
            "iRidi SCADA-BMS  ·  Stimulsoft Reports Engine  ·  devbms.iridi.com",
            font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# =========================================================
# СОХРАНЕНИЕ
# =========================================================
out_path = "iRidi_BMS_Reporting_System.pptx"
prs.save(out_path)
print(f"✓ Презентация сохранена: {out_path}")
print(f"✓ Слайдов: {len(prs.slides)}")