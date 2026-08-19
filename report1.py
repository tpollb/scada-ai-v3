# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Цветовая схема — тёмная тема с акцентами
BG_DARK = RGBColor(18, 18, 18)           # #121212
BG_PANEL = RGBColor(28, 28, 28)          # #1C1C1C
BG_ACCENT = RGBColor(35, 35, 42)         # #23232A
TEXT_PRIMARY = RGBColor(240, 240, 240)   # #F0F0F0
TEXT_SECONDARY = RGBColor(170, 170, 170) # #AAAAAA
TEXT_DIM = RGBColor(120, 120, 120)       # #787878
ACCENT_BLUE = RGBColor(0, 150, 255)      # #0096FF
ACCENT_CYAN = RGBColor(0, 200, 200)      # #00C8C8
ACCENT_GREEN = RGBColor(80, 200, 120)    # #50C878
ACCENT_AMBER = RGBColor(255, 180, 0)     # #FFB400
ACCENT_PURPLE = RGBColor(160, 120, 255)  # #A078FF
LINE_COLOR = RGBColor(55, 55, 55)        # #373737

FONT_LIGHT = 'Segoe UI Light'
FONT_REGULAR = 'Segoe UI'
FONT_SEMIBOLD = 'Segoe UI Semibold'

def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text,
             font_size=14, color=TEXT_PRIMARY, bold=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT_LIGHT,
             italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, default_size=13,
                  default_color=TEXT_PRIMARY, default_font=FONT_LIGHT):
    """lines = список кортежей (text, size, color, bold, italic, space_after)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text = line[0]
        size = line[1] if len(line) > 1 else default_size
        color = line[2] if len(line) > 2 else default_color
        bold = line[3] if len(line) > 3 else False
        italic = line[4] if len(line) > 4 else False
        space_after = line[5] if len(line) > 5 else 6

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = default_font
        p.font.italic = italic
        p.space_after = Pt(space_after)
    return tf

def add_panel(slide, left, top, width, height, color=BG_PANEL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, height, color=ACCENT_BLUE, width=0.05):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color=LINE_COLOR, thickness=1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_header(slide, tag, title):
    """Стандартный заголовок слайда: тег + заголовок + линия"""
    add_text(slide, 1.0, 0.5, 11.0, 0.4, tag,
             font_size=11, color=ACCENT_BLUE, font_name=FONT_SEMIBOLD)
    add_text(slide, 1.0, 0.85, 11.0, 0.7, title,
             font_size=30, color=TEXT_PRIMARY, font_name=FONT_LIGHT)
    add_line(slide, 1.0, 1.6, 2.2, ACCENT_BLUE, 3)

# =====================================================================
# СЛАЙД 1: ТИТУЛЬНЫЙ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_line(slide, 1.0, 1.6, 3.5, ACCENT_BLUE, 3)

add_text(slide, 1.0, 1.9, 11.0, 1.3,
         "СИСТЕМА ОТЧЁТОВ\niRidi SCADA-BMS",
         font_size=52, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

add_text(slide, 1.0, 3.4, 11.0, 0.9,
         "Ценность для бизнеса через технологии отчётности",
         font_size=24, color=ACCENT_BLUE, font_name=FONT_LIGHT)

add_line(slide, 1.0, 4.5, 11.0, LINE_COLOR, 1)

add_text(slide, 1.0, 4.8, 9.0, 1.2,
         "От инженерных данных — к управленческим решениям,\n"
         "финансовым документам и юридической защите бизнеса.",
         font_size=17, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

add_panel(slide, 1.0, 6.2, 11.3, 0.9, BG_PANEL)
add_text(slide, 1.4, 6.35, 10.5, 0.6,
         "Платформы: Windows · Linux · Debian 12 · Сервер кластера    |    "
         "Движок: Stimulsoft Reports    |    "
         "Источники: PostgreSQL · OData · REST API",
         font_size=12, color=TEXT_DIM, font_name=FONT_LIGHT)

# =====================================================================
# СЛАЙД 2: ЧТО ЭТО
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "01 · ЧТО ЭТО", "Система отчётов iRidi SCADA-BMS")

add_text(slide, 1.0, 2.0, 11.3, 1.2,
         "Встроенный модуль формирования, визуализации и экспорта отчётности "
         "по всем инженерным системам объекта — от отдельного датчика "
         "до распределённой сети зданий.",
         font_size=18, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

# Три блока — что включает
blocks = [
    ("ГОТОВЫЕ ШАБЛОНЫ", ACCENT_BLUE,
     ["История событий", "Значения тегов", "Аварии и подтверждения",
      "Действия персонала", "Состояние БД и системы"]),
    ("ИНСТРУМЕНТЫ НАСТРОЙКИ", ACCENT_CYAN,
     ["Веб-дизайнер отчётов", "Группировка по тегам и датам",
      "Фильтры и сортировки", "Итоги и агрегации",
      "Темы и форматы оформления"]),
    ("ФОРМАТЫ ВЫДАЧИ", ACCENT_GREEN,
     ["PDF с подписями и PDF/A", "Excel, Word, CSV",
      "HTML, JSON, XML, SVG", "Печать напрямую из браузера",
      "Экспорт по расписанию"]),
]

for i, (title, color, items) in enumerate(blocks):
    x = 1.0 + i * 3.95
    add_panel(slide, x, 3.4, 3.7, 3.5, BG_PANEL)
    add_accent_bar(slide, x, 3.4, 3.5, color)
    add_text(slide, x + 0.25, 3.6, 3.3, 0.4, title,
             font_size=12, color=color, font_name=FONT_SEMIBOLD)
    lines = [(f"›  {item}", 12, TEXT_PRIMARY, False, False, 5) for item in items]
    add_multiline(slide, x + 0.25, 4.1, 3.3, 2.6, lines)

# =====================================================================
# СЛАЙД 3: ЦКП — ГЛАВНЫЙ СЛАЙД
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "02 · ЦКП", "Ценный конечный продукт системы отчётности")

add_panel(slide, 1.0, 2.0, 11.3, 1.6, BG_ACCENT)
add_accent_bar(slide, 1.0, 2.0, 1.6, ACCENT_BLUE)
add_text(slide, 1.3, 2.2, 10.8, 0.4, "ЦКП",
         font_size=12, color=ACCENT_BLUE, font_name=FONT_SEMIBOLD)
add_text(slide, 1.3, 2.55, 10.8, 0.9,
         "Превращение сырых инженерных данных в управленческие решения, "
         "финансовые документы и юридически значимую доказательную базу — "
         "без привлечения программистов и внешних BI-систем.",
         font_size=20, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

# Четыре ценности
values = [
    ("01", "ПРОЗРАЧНОСТЬ", ACCENT_BLUE,
     "Полная картина работы инженерных систем "
     "с точностью до секунды — от состояния "
     "датчика до действия оператора."),
    ("02", "ЗАЩИТА", ACCENT_AMBER,
     "Юридически значимые документы для споров "
     "с ресурсоснабжающими организациями, "
     "аудитов и проверок."),
    ("03", "ЭКОНОМИКА", ACCENT_GREEN,
     "Снижение затрат на энергоресурсы, "
     "сокращение трудозатрат на отчётность "
     "и возврат переплат через доказательную базу."),
    ("04", "СКОРОСТЬ", ACCENT_CYAN,
     "Отчёты формируются за секунды вместо часов. "
     "Шаблон настраивается один раз — "
     "далее работает автоматически."),
]

for i, (num, title, color, desc) in enumerate(values):
    x = 1.0 + i * 2.95
    add_panel(slide, x, 4.0, 2.75, 2.9, BG_PANEL)
    add_accent_bar(slide, x, 4.0, 2.9, color)
    add_text(slide, x + 0.25, 4.15, 0.5, 0.4, num,
             font_size=14, color=color, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.25, 4.5, 2.3, 0.4, title,
             font_size=13, color=TEXT_PRIMARY, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.25, 4.9, 2.3, 1.8, desc,
             font_size=12, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

# =====================================================================
# СЛАЙД 4: ЦЕННОСТЬ 1 — ПРОЗРАЧНОСТЬ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "03 · ЦЕННОСТЬ", "Прозрачность эксплуатации")

# Левая часть — ценность
add_panel(slide, 1.0, 2.0, 5.6, 4.9, BG_PANEL)
add_accent_bar(slide, 1.0, 2.0, 4.9, ACCENT_BLUE)
add_text(slide, 1.3, 2.2, 5.1, 0.4, "ЧТО ПОЛУЧАЕТ ЗАКАЗЧИК",
         font_size=12, color=ACCENT_BLUE, font_name=FONT_SEMIBOLD)

add_multiline(slide, 1.3, 2.7, 5.1, 3.8, [
    ("Полная наблюдаемость инженерных систем", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Каждое событие фиксируется с точностью до секунды", 12, TEXT_SECONDARY, False, False, 6),
    ("›  19 типов событий: аварии, действия персонала, "
     "изменения тегов, состояния БД", 12, TEXT_SECONDARY, False, False, 6),
    ("›  История доступна за любой период: час, день, "
     "месяц, год", 12, TEXT_SECONDARY, False, False, 12),
    ("Контроль действий персонала", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Фиксация входов и выходов пользователей", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Логирование всех изменений параметров", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Подтверждение и квитирование аварий", 12, TEXT_SECONDARY, False, False, 12),
    ("Единая точка правды для всех служб", 15, TEXT_PRIMARY, True, False, 8),
    ("Диспетчер, инженер, руководство и аудит работают "
     "с одним набором данных", 12, TEXT_SECONDARY, False, False, 6),
])

# Правая часть — за счёт чего
add_panel(slide, 6.9, 2.0, 5.4, 4.9, BG_PANEL)
add_accent_bar(slide, 6.9, 2.0, 4.9, ACCENT_CYAN)
add_text(slide, 7.2, 2.2, 4.9, 0.4, "ЗА СЧЁТ ЧЕГО ДОСТИГАЕТСЯ",
         font_size=12, color=ACCENT_CYAN, font_name=FONT_SEMIBOLD)

add_multiline(slide, 7.2, 2.7, 4.9, 3.8, [
    ("Технологический фундамент", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Движок Stimulsoft Reports — промышленный "
     "стандарт отчётности", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Подключение напрямую к PostgreSQL SCADA-сервера", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Серверная генерация отчётов — не нагрузка "
     "на клиентские рабочие места", 12, TEXT_SECONDARY, False, False, 12),
    ("Гибкая аналитика", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Группировка по тегам, зонам, датам "
     "(день / неделя / месяц)", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Функции агрегации: Sum, Count, Average, "
     "DayOfYear и другие", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Фильтрация по любым полям, "
     "включая не включённые в отчёт", 12, TEXT_SECONDARY, False, False, 12),
    ("Мгновенный доступ", 15, TEXT_PRIMARY, True, False, 8),
    ("Веб-интерфейс из любого браузера — без установки "
     "дополнительного ПО на рабочем месте", 12, TEXT_SECONDARY, False, False, 6),
])

# =====================================================================
# СЛАЙД 5: ЦЕННОСТЬ 2 — ЮРИДИЧЕСКАЯ ЗАЩИТА
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "04 · ЦЕННОСТЬ", "Юридическая защита и доказательная база")

add_panel(slide, 1.0, 2.0, 5.6, 4.9, BG_PANEL)
add_accent_bar(slide, 1.0, 2.0, 4.9, ACCENT_AMBER)
add_text(slide, 1.3, 2.2, 5.1, 0.4, "ЧТО ПОЛУЧАЕТ ЗАКАЗЧИК",
         font_size=12, color=ACCENT_AMBER, font_name=FONT_SEMIBOLD)

add_multiline(slide, 1.3, 2.7, 5.1, 3.8, [
    ("Защита в спорах с РСО", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Доказательная база при расхождениях "
     "в показаниях с энергосбытом, водоканалом, "
     "теплосетью", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Подтверждение периодов простоя "
     "и отсутствия потребления", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Возврат переплат через приложения "
     "к претензиям и исковым заявлениям", 12, TEXT_SECONDARY, False, False, 12),
    ("Защита при аудитах и проверках", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Подтверждение соблюдения регламентов "
     "эксплуатации", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Документирование реакции персонала "
     "на аварийные ситуации", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Доказательство корректной работы "
     "инженерных систем", 12, TEXT_SECONDARY, False, False, 12),
    ("Документооборот с бухгалтерией", 15, TEXT_PRIMARY, True, False, 8),
    ("Готовые выгрузки для сверки с ресурсоснабжающими "
     "организациями и внутренней отчётности", 12, TEXT_SECONDARY, False, False, 6),
])

add_panel(slide, 6.9, 2.0, 5.4, 4.9, BG_PANEL)
add_accent_bar(slide, 6.9, 2.0, 4.9, ACCENT_PURPLE)
add_text(slide, 7.2, 2.2, 4.9, 0.4, "ЗА СЧЁТ ЧЕГО ДОСТИГАЕТСЯ",
         font_size=12, color=ACCENT_PURPLE, font_name=FONT_SEMIBOLD)

add_multiline(slide, 7.2, 2.7, 4.9, 3.8, [
    ("Неизменяемость и полнота данных", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Данные хранятся в PostgreSQL с привязкой "
     "к временным меткам", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Логирование 19 типов событий "
     "исключает пробелы в истории", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Фиксация действий каждого пользователя "
     "персонально", 12, TEXT_SECONDARY, False, False, 12),
    ("Юридически значимые форматы", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Экспорт в PDF с поддержкой цифровых "
     "подписей и шифрования", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Формат PDF/A для долгосрочного "
     "архивного хранения", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Экспорт в Excel для построчной сверки "
     "с актами РСО", 12, TEXT_SECONDARY, False, False, 12),
    ("Консолидация данных", 15, TEXT_PRIMARY, True, False, 8),
    ("Группировка по расчётным периодам "
     "(день / неделя / месяц) формирует готовую "
     "картину для юристов и бухгалтерии", 12, TEXT_SECONDARY, False, False, 6),
])

# =====================================================================
# СЛАЙД 6: ЦЕННОСТЬ 3 — ЭКОНОМИКА
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "05 · ЦЕННОСТЬ", "Экономическая эффективность")

# Четыре карточки с цифрами
cards = [
    ("15–25%", "снижение затрат\nна энергоресурсы",
     "Выявление аномалий потребления "
     "и неэффективных режимов работы "
     "оборудования через анализ "
     "значений тегов за периоды.",
     ACCENT_GREEN),
    ("30–50%", "сокращение трудозатрат\nна отчётность",
     "Автоматические шаблоны "
     "заменяют ручной сбор данных "
     "из журналов. Отчёт, который "
     "делался часами, — за секунды.",
     ACCENT_BLUE),
    ("100%", "возврат переплат\nпри оспаривании",
     "При успешном оспаривании "
     "начислений РСО — полный возврат "
     "на основании доказательной "
     "базы из SCADA.",
     ACCENT_AMBER),
    ("0 ₽", "дополнительных\nотчислений",
     "Движок Stimulsoft лицензирован "
     "royalty-free. Создание отчётов "
     "конечными пользователями не "
     "требует дополнительных платежей.",
     ACCENT_CYAN),
]

for i, (num, title, desc, color) in enumerate(cards):
    x = 1.0 + i * 2.95
    add_panel(slide, x, 2.0, 2.75, 4.9, BG_PANEL)
    add_accent_bar(slide, x, 2.0, 4.9, color)
    add_text(slide, x + 0.25, 2.2, 2.3, 0.7, num,
             font_size=36, color=color, font_name=FONT_LIGHT)
    add_text(slide, x + 0.25, 2.95, 2.3, 0.7, title,
             font_size=12, color=TEXT_PRIMARY, font_name=FONT_SEMIBOLD)
    add_line(slide, x + 0.25, 3.75, 2.2, LINE_COLOR, 1)
    add_text(slide, x + 0.25, 3.9, 2.3, 2.8, desc,
             font_size=11, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

# Нижний блок — итоговая окупаемость
add_panel(slide, 1.0, 7.05, 11.3, 0.35, BG_ACCENT)
add_text(slide, 1.3, 7.1, 10.8, 0.25,
         "Типовая окупаемость модуля: 3–6 месяцев только за счёт оспаривания начислений "
         "и оптимизации энергопотребления.",
         font_size=12, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

# =====================================================================
# СЛАЙД 7: ЦЕННОСТЬ 4 — СКОРОСТЬ И ИНТЕГРАЦИЯ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "06 · ЦЕННОСТЬ", "Скорость и интеграция в бизнес-процессы")

add_panel(slide, 1.0, 2.0, 5.6, 4.9, BG_PANEL)
add_accent_bar(slide, 1.0, 2.0, 4.9, ACCENT_CYAN)
add_text(slide, 1.3, 2.2, 5.1, 0.4, "ЧТО ПОЛУЧАЕТ ЗАКАЗЧИК",
         font_size=12, color=ACCENT_CYAN, font_name=FONT_SEMIBOLD)

add_multiline(slide, 1.3, 2.7, 5.1, 3.8, [
    ("Мгновенная готовность отчётов", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Шаблон настраивается один раз — "
     "далее работает автоматически", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Формирование любого отчёта — секунды, "
     "не часы ручного труда", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Регулярные регламентные отчёты — "
     "без участия оператора", 12, TEXT_SECONDARY, False, False, 12),
    ("Интеграция в документооборот", 15, TEXT_PRIMARY, True, False, 12),
    ("›  Бухгалтерия получает выгрузки "
     "в нужных форматах (Excel, CSV, PDF)", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Руководство — сводки в удобном "
     "для принятия решений виде", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Сервисные службы — отчёты "
     "по инцидентам для планирования ППР", 12, TEXT_SECONDARY, False, False, 12),
    ("Масштабируемость решения", 15, TEXT_PRIMARY, True, False, 8),
    ("От одного здания до распределённой "
     "сети объектов — единая система "
     "отчётности без перенастройки", 12, TEXT_SECONDARY, False, False, 6),
])

add_panel(slide, 6.9, 2.0, 5.4, 4.9, BG_PANEL)
add_accent_bar(slide, 6.9, 2.0, 4.9, ACCENT_GREEN)
add_text(slide, 7.2, 2.2, 4.9, 0.4, "ЗА СЧЁТ ЧЕГО ДОСТИГАЕТСЯ",
         font_size=12, color=ACCENT_GREEN, font_name=FONT_SEMIBOLD)

add_multiline(slide, 7.2, 2.7, 4.9, 3.8, [
    ("Широкий спектр форматов экспорта", 15, TEXT_PRIMARY, True, False, 12),
    ("›  PDF с подписями и PDF/A для архивов", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Excel и CSV для бухгалтерии и аналитики", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Word и RTF для текстовых документов", 12, TEXT_SECONDARY, False, False, 6),
    ("›  HTML для веб-представления", 12, TEXT_SECONDARY, False, False, 6),
    ("›  JSON, XML для интеграции с ERP/BI", 12, TEXT_SECONDARY, False, False, 12),
    ("Гибкие источники данных", 15, TEXT_PRIMARY, True, False, 12),
    ("›  PostgreSQL, OData, REST API", 12, TEXT_SECONDARY, False, False, 6),
    ("›  JSON, XML, Excel как файловые источники", 12, TEXT_SECONDARY, False, False, 6),
    ("›  Поддержка SQL и NoSQL баз данных", 12, TEXT_SECONDARY, False, False, 12),
    ("Кроссплатформенность", 15, TEXT_PRIMARY, True, False, 8),
    ("Windows, Linux, Debian 12, Сервер кластера — "
     "единый движок Stimulsoft на всех платформах", 12, TEXT_SECONDARY, False, False, 6),
])

# =====================================================================
# СЛАЙД 8: ТЕХНОЛОГИЧЕСКАЯ БАЗА
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "07 · ТЕХНОЛОГИИ", "За счёт чего достигается ценность")

add_text(slide, 1.0, 2.0, 11.3, 0.9,
         "Четыре технологических слоя, которые превращают инженерные "
         "данные в бизнес-результат.",
         font_size=16, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

layers = [
    ("СЛОЙ 1 · ДВИЖОК ОТЧЁТНОСТИ", ACCENT_BLUE,
     "Stimulsoft Reports",
     ["Промышленный стандарт отчётности",
      "Чистый JavaScript — работа в любом современном браузере",
      "Поддержка Node.js для серверной генерации",
      "Лицензия royalty-free — без отчислений за пользователей",
      "Встроенный ИИ-ассистент в дизайнере"]),
    ("СЛОЙ 2 · ИСТОЧНИКИ ДАННЫХ", ACCENT_CYAN,
     "PostgreSQL + интеграции",
     ["Прямое подключение к БД SCADA-сервера",
      "MS SQL, MySQL, Oracle, Firebird, MongoDB",
      "OData, REST API, JSON, XML, Excel",
      "ODBC и PDO — подключение к любым СУБД",
      "Удалённое расположение БД без потери производительности"]),
    ("СЛОЙ 3 · ИНСТРУМЕНТЫ АНАЛИТИКИ", ACCENT_GREEN,
     "Дизайнер и вьювер",
     ["Веб-дизайнер отчётов прямо из браузера",
      "Группировка, фильтрация, сортировка, итоги",
      "Master-Detail отчёты и многоуровневая детализация",
      "Функции агрегации: Sum, Count, Average, DayOfYear",
      "Настраиваемый вьювер с темами оформления"]),
    ("СЛОЙ 4 · ФОРМАТЫ ВЫДАЧИ", ACCENT_AMBER,
     "25+ форматов экспорта",
     ["PDF с поддержкой подписей и шифрования",
      "PDF/A для долгосрочного архивного хранения",
      "Excel, Word, CSV для бизнес-пользователей",
      "HTML, JSON, XML для интеграции с ERP/BI",
      "Печать и экспорт на стороне клиента"]),
]

for i, (title, color, subtitle, items) in enumerate(layers):
    x = 1.0 + i * 2.95
    add_panel(slide, x, 3.1, 2.75, 4.15, BG_PANEL)
    add_accent_bar(slide, x, 3.1, 4.15, color)
    add_text(slide, x + 0.2, 3.25, 2.35, 0.35, title,
             font_size=10, color=color, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.2, 3.6, 2.35, 0.35, subtitle,
             font_size=12, color=TEXT_PRIMARY, font_name=FONT_SEMIBOLD)
    lines = [(f"›  {item}", 10, TEXT_SECONDARY, False, False, 4) for item in items]
    add_multiline(slide, x + 0.2, 4.05, 2.35, 3.1, lines)

# =====================================================================
# СЛАЙД 9: АРХИТЕКТУРА ЦЕННОСТИ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "08 · АРХИТЕКТУРА", "Как данные превращаются в бизнес-решения")

# Горизонтальная цепочка преобразования
steps = [
    ("ДАННЫЕ", ACCENT_BLUE,
     "19 типов событий\nот SCADA-сервера",
     "Аварии, действия персонала,\nизменения тегов, логины,\nсостояния БД"),
    ("ХРАНЕНИЕ", ACCENT_CYAN,
     "PostgreSQL\nс временными метками",
     "Неизменяемый лог событий\nс точностью до секунды\nи привязкой к зонам"),
    ("ОБРАБОТКА", ACCENT_GREEN,
     "Stimulsoft Reports\n+ шаблоны",
     "Группировка, фильтрация,\nагрегация, расчёт итогов\nпо заданным правилам"),
    ("ВЫДАЧА", ACCENT_AMBER,
     "PDF · Excel · CSV\n· JSON · HTML",
     "Готовые документы\nдля бизнеса, бухгалтерии,\nюристов и руководства"),
    ("РЕШЕНИЕ", ACCENT_PURPLE,
     "Управленческое\nдействие",
     "Оспаривание начислений,\nоптимизация энергопотребления,\nпланирование ППР"),
]

for i, (title, color, subtitle, desc) in enumerate(steps):
    x = 0.6 + i * 2.55
    add_panel(slide, x, 2.0, 2.35, 3.5, BG_PANEL)
    add_accent_bar(slide, x, 2.0, 3.5, color)
    add_text(slide, x + 0.2, 2.2, 2.0, 0.35, f"0{i+1}",
             font_size=11, color=color, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.2, 2.5, 2.0, 0.35, title,
             font_size=13, color=TEXT_PRIMARY, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.2, 2.9, 2.0, 0.6, subtitle,
             font_size=11, color=color, font_name=FONT_LIGHT)
    add_line(slide, x + 0.2, 3.6, 1.95, LINE_COLOR, 1)
    add_text(slide, x + 0.2, 3.75, 2.0, 1.5, desc,
             font_size=10, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

# Нижний блок — итоговая формула ценности
add_panel(slide, 0.6, 5.8, 12.15, 1.4, BG_ACCENT)
add_accent_bar(slide, 0.6, 5.8, 1.4, ACCENT_BLUE)
add_text(slide, 0.9, 5.95, 11.6, 0.35, "ФОРМУЛА ЦЕННОСТИ",
         font_size=12, color=ACCENT_BLUE, font_name=FONT_SEMIBOLD)
add_text(slide, 0.9, 6.35, 11.6, 0.7,
         "Инженерные данные  +  промышленный движок отчётности  +  "
         "гибкие инструменты аналитики  +  юридически значимые форматы\n"
         "=  управленческие решения, финансовая экономия и юридическая защита бизнеса.",
         font_size=14, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

# =====================================================================
# СЛАЙД 10: СЦЕНАРИИ ПРИМЕНЕНИЯ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "09 · СЦЕНАРИИ", "Где система приносит измеримую ценность")

scenarios = [
    ("СПОРЫ С РЕСУРСОСНАБЖАЮЩИМИ ОРГАНИЗАЦИЯМИ",
     ACCENT_AMBER,
     "Главный инженер, юрист, бухгалтерия",
     "Формирование отчётов по показаниям приборов учёта "
     "за расчётный период с группировкой по дням. "
     "Экспорт в PDF и Excel как приложение к претензии.",
     "Возврат 100% переплат при успешном оспаривании"),
    ("ЭНЕРГЕТИЧЕСКИЙ АУДИТ",
     ACCENT_GREEN,
     "Энергоменеджер, главный энергетик",
     "Анализ значений тегов счётчиков и датчиков "
     "с группировкой по периодам. Выявление пиков "
     "и аномалий потребления.",
     "Снижение затрат на энергоресурсы на 15–25%"),
    ("РЕГЛАМЕНТНАЯ ОТЧЁТНОСТЬ",
     ACCENT_BLUE,
     "Руководство, собственник, арендодатель",
     "Автоматические ежемесячные сводки по KPI "
     "эксплуатации, статистике отказов, потреблению "
     "ресурсов без участия оператора.",
     "Сокращение трудозатрат на отчётность на 30–50%"),
    ("АУДИТ ДЕЙСТВИЙ ПЕРСОНАЛА",
     ACCENT_PURPLE,
     "Служба безопасности, руководство объекта",
     "Отчёты по входам, выходам и действиям "
     "операторов. Фиксация несанкционированного "
     "доступа и изменений уставок.",
     "Доказательная база при расследовании инцидентов"),
    ("АВАРИЙНЫЙ АНАЛИЗ",
     ACCENT_CYAN,
     "Служба эксплуатации, сервисная служба",
     "Последовательность событий при инцидентах, "
     "время реакции персонала, длительность "
     "аварий. Выявление системных проблем.",
     "Снижение повторяемости инцидентов"),
    ("СВЕРКА С БУХГАЛТЕРИЕЙ",
     ACCENT_AMBER,
     "Бухгалтерия, финансовый отдел",
     "Построчная сверка данных SCADA с актами "
     "ресурсоснабжающих организаций. Выявление "
     "расхождений до оплаты.",
     "Предотвращение необоснованных платежей"),
]

for i, (title, color, who, desc, value) in enumerate(scenarios):
    row = i // 3
    col = i % 3
    x = 1.0 + col * 3.95
    y = 2.0 + row * 2.75
    add_panel(slide, x, y, 3.7, 2.55, BG_PANEL)
    add_accent_bar(slide, x, y, 2.55, color)
    add_text(slide, x + 0.25, y + 0.15, 3.3, 0.3, title,
             font_size=10, color=color, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.25, y + 0.45, 3.3, 0.2, who,
             font_size=9, color=TEXT_DIM, font_name=FONT_LIGHT, italic=True)
    add_text(slide, x + 0.25, y + 0.75, 3.3, 1.1, desc,
             font_size=10, color=TEXT_PRIMARY, font_name=FONT_LIGHT)
    add_line(slide, x + 0.25, y + 1.95, 3.2, LINE_COLOR, 1)
    add_text(slide, x + 0.25, y + 2.05, 3.3, 0.4,
             f"Ценность: {value}",
             font_size=10, color=color, font_name=FONT_SEMIBOLD)

# =====================================================================
# СЛАЙД 11: ИТОГИ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10 · ИТОГИ", "Что получает заказчик")

add_text(slide, 1.0, 2.0, 11.3, 0.8,
         "Система отчётов iRidi SCADA-BMS — это не просто модуль "
         "визуализации данных. Это инструмент превращения инженерной "
         "информации в бизнес-результат.",
         font_size=17, color=TEXT_PRIMARY, font_name=FONT_LIGHT)

results = [
    ("01", "ПРОЗРАЧНОСТЬ", ACCENT_BLUE,
     "Полная наблюдаемость инженерных систем "
     "с точностью до секунды. Контроль действий "
     "персонала. Единая точка правды для всех служб."),
    ("02", "ЮРИДИЧЕСКАЯ ЗАЩИТА", ACCENT_AMBER,
     "Доказательная база для споров с РСО, "
     "аудитов и проверок. Юридически значимые "
     "форматы с поддержкой подписей и PDF/A."),
    ("03", "ЭКОНОМИЯ", ACCENT_GREEN,
     "Снижение затрат на энергоресурсы на 15–25%. "
     "Сокращение трудозатрат на отчётность на 30–50%. "
     "Возврат переплат при оспаривании."),
    ("04", "СКОРОСТЬ", ACCENT_CYAN,
     "Отчёты за секунды вместо часов. "
     "Шаблон настраивается один раз. "
     "Регламентные отчёты формируются автоматически."),
    ("05", "ИНТЕГРАЦИЯ", ACCENT_PURPLE,
     "25+ форматов экспорта. Готовые выгрузки "
     "для бухгалтерии, юристов, руководства. "
     "Интеграция с ERP и BI через JSON и XML."),
    ("06", "МАСШТАБИРУЕМОСТЬ", ACCENT_BLUE,
     "От одного здания до распределённой сети "
     "объектов. Кроссплатформенность: "
     "Windows, Linux, Debian 12, Сервер кластера."),
]

for i, (num, title, color, desc) in enumerate(results):
    row = i // 3
    col = i % 3
    x = 1.0 + col * 3.95
    y = 3.0 + row * 2.1
    add_panel(slide, x, y, 3.7, 1.9, BG_PANEL)
    add_accent_bar(slide, x, y, 1.9, color)
    add_text(slide, x + 0.25, y + 0.1, 0.5, 0.3, num,
             font_size=13, color=color, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.25, y + 0.4, 3.3, 0.3, title,
             font_size=12, color=TEXT_PRIMARY, font_name=FONT_SEMIBOLD)
    add_text(slide, x + 0.25, y + 0.75, 3.3, 1.0, desc,
             font_size=11, color=TEXT_SECONDARY, font_name=FONT_LIGHT)

# =====================================================================
# СОХРАНЕНИЕ
# =====================================================================
out_path = 'iRidi_Reports_Ценность_для_бизнеса.pptx'
prs.save(out_path)
print(f"✓ Презентация сохранена: {out_path}")
print(f"✓ Слайдов: {len(prs.slides)}")
print(f"✓ Размер: {len(open(out_path, 'rb').read())/1024:.1f} KB")