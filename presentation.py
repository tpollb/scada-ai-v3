#!/usr/bin/env python3
"""
SCADA.AI Ценовая стратегия
Внутренняя презентация для руководства АО ИРИДИЙ БМС
Реалистичные прогнозы | Строгий деловой стиль | Segoe UI Light
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
matplotlib.rcParams['font.family'] = 'Segoe UI'

# === Строгая деловая палитра ===
BG_COLOR = RGBColor(248, 250, 252)        # slate-50
BG_DARK = RGBColor(15, 23, 42)           # slate-900
BG_CARD = RGBColor(241, 245, 249)        # slate-100
TEXT_PRIMARY = RGBColor(15, 23, 42)      # slate-900
TEXT_SECONDARY = RGBColor(71, 85, 105)   # slate-600
TEXT_LIGHT = RGBColor(241, 245, 249)     # slate-100
ACCENT_NAVY = RGBColor(30, 64, 175)      # blue-800
ACCENT_TEAL = RGBColor(15, 118, 110)     # teal-700
ACCENT_AMBER = RGBColor(180, 83, 9)      # amber-700
ACCENT_RED = RGBColor(153, 27, 27)       # red-800
ACCENT_GREEN = RGBColor(21, 128, 61)     # green-700
DEEP_PURPLE = RGBColor(88, 28, 135)
DIVIDER = RGBColor(203, 213, 225)        # slate-300

FONT_NAME = 'Segoe UI Light'
FONT_NAME_BOLD = 'Segoe UI Semibold'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_font(run, name=FONT_NAME, size=18, bold=False, color=TEXT_PRIMARY):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color=BG_CARD, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=TEXT_PRIMARY,
             bold=False, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=size, bold=bold, color=color)
    p.alignment = align
    return txBox

def add_title_bar(slide, title, subtitle=""):
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.55), title,
             size=26, bold=False, color=TEXT_LIGHT, font_name=FONT_NAME)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.25), subtitle,
                 size=11, color=RGBColor(148, 163, 184), font_name=FONT_NAME)

def add_footer(slide, text="АО ИРИДИЙ БМС  |  Июль 2026  |  SCADA.AI v3.2.9.1"):
    add_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3),
             text, size=9, color=TEXT_SECONDARY, font_name=FONT_NAME)

# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1),
         "SCADA.AI", size=80, bold=False, color=TEXT_LIGHT,
         align=PP_ALIGN.CENTER, font_name=FONT_NAME)

add_shape_bg(slide, Inches(5.5), Inches(3.3), Inches(2.333), Inches(0.02), ACCENT_NAVY)

add_text(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.7),
         "Ценовая стратегия и модель монетизации",
         size=28, bold=False, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

add_shape_bg(slide, Inches(3), Inches(4.8), Inches(7.333), Inches(0.9), RGBColor(30, 41, 59))
add_text(slide, Inches(3.3), Inches(4.95), Inches(6.733), Inches(0.6),
         "Внутренний документ для обсуждения",
         size=18, bold=False, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

add_text(slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.4),
         "АО ИРИДИЙ БМС   |   Июль 2026   |   Версия 3.2.9.1",
         size=12, color=RGBColor(148, 163, 184), align=PP_ALIGN.CENTER, font_name=FONT_NAME)

# ============================================================
# SLIDE 2: Три тарифа — обзор
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Три тарифных плана", "Сегментация по потребностям и бюджетам клиентов")

tiers = [
    ("STARTER", "Базовое здоровье здания",
     "2 500 — 5 000 тегов",
     "60 000 — 90 000 ₽/мес",
     [
         "Мониторинг параметров среды",
         "Статус оборудования",
         "Базовые алерты",
         "Тренды и экспорт данных",
     ], ACCENT_TEAL),
    ("PROFESSIONAL", "Аналитика и диагностика",
     "15 000 — 50 000 тегов",
     "300 000 — 750 000 ₽/мес",
     [
         "Всё из Starter",
         "Глубокая диагностика (DDA)",
         "A/B анализ и сезонность",
         "LLM-интерпретация",
     ], ACCENT_NAVY),
    ("ENTERPRISE", "Максимальная ценность",
     "150 000+ тегов",
     "от 1 800 000 ₽/мес",
     [
         "Всё из Professional",
         "On-premise LLM",
         "Custom ML модели",
         "Мульти-объекты и SLA",
     ], DEEP_PURPLE),
]

for i, (name, desc, tags, price, features, color) in enumerate(tiers):
    left = Inches(0.4 + i * 4.25)
    add_shape_bg(slide, left, Inches(1.6), Inches(4.0), Inches(5.3), BG_CARD, line_color=DIVIDER)
    add_shape_bg(slide, left, Inches(1.6), Inches(4.0), Inches(0.08), color)

    add_text(slide, left, Inches(1.85), Inches(4.0), Inches(0.4), name,
             size=20, bold=True, color=color, align=PP_ALIGN.CENTER, font_name=FONT_NAME_BOLD)
    add_text(slide, left + Inches(0.3), Inches(2.3), Inches(3.4), Inches(0.4), desc,
             size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

    add_shape_bg(slide, left + Inches(0.3), Inches(2.8), Inches(3.4), Inches(0.005), DIVIDER)

    add_text(slide, left + Inches(0.3), Inches(2.95), Inches(3.4), Inches(0.35), tags,
             size=11, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

    add_shape_bg(slide, left + Inches(0.3), Inches(3.4), Inches(3.4), Inches(0.55), BG_DARK)
    add_text(slide, left + Inches(0.3), Inches(3.48), Inches(3.4), Inches(0.4), price,
             size=14, bold=False, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

    for j, feat in enumerate(features):
        add_text(slide, left + Inches(0.4), Inches(4.2 + j * 0.42), Inches(3.3), Inches(0.4),
                 f"—  {feat}", size=11, color=TEXT_PRIMARY, font_name=FONT_NAME)

add_footer(slide)

# ============================================================
# SLIDE 3: Starter детально
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Тариф Starter", "Базовое решение для малого и среднего сегмента")

# Левая часть
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(0.08), Inches(5.3), ACCENT_TEAL)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.4), "Включено в тариф",
         size=16, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

features_starter = [
    "Health Score здания (комплексная оценка 0-100)",
    "Мониторинг 5 параметров: CO2, VOC, температура, влажность, давление",
    "Статус оборудования: online / offline / stuck",
    "Простые алерты по email и telegram",
    "Тренды за 7 / 30 / 90 дней",
    "Экспорт данных в CSV и Excel",
    "Базовый дашборд с ключевыми метриками",
    "1 пользователь с правами администратора",
]
for i, feat in enumerate(features_starter):
    add_text(slide, Inches(1.0), Inches(2.4 + i * 0.47), Inches(5.5), Inches(0.4),
             f"—  {feat}", size=11, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Правая часть — цены
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.833), Inches(2.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(0.08), Inches(2.3), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(1.8), Inches(5), Inches(0.4), "Ценообразование",
         size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

add_text(slide, Inches(7.4), Inches(2.3), Inches(5), Inches(1.4),
         "2 500 тегов:  60 000 ₽/мес  |  612 000 ₽/год\n"
         "5 000 тегов:  90 000 ₽/мес  |  918 000 ₽/год\n\n"
         "Маржа: 20-25%   |   Окупаемость для клиента: 6-8 мес",
         size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Правая часть — целевые клиенты
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(5.833), Inches(2.8), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(0.08), Inches(2.8), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(4.3), Inches(5), Inches(0.4), "Целевые сегменты",
         size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

customers = [
    ("Образовательные учреждения", "40 000 — 60 000 ₽/мес"),
    ("Медицинские учреждения", "60 000 — 90 000 ₽/мес"),
    ("Офисные здания", "50 000 — 70 000 ₽/мес"),
    ("Торговые центры (малые)", "70 000 — 90 000 ₽/мес"),
]
for i, (name, price) in enumerate(customers):
    add_text(slide, Inches(7.4), Inches(4.9 + i * 0.45), Inches(3.2), Inches(0.4),
             name, size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)
    add_text(slide, Inches(10.3), Inches(4.9 + i * 0.45), Inches(2.2), Inches(0.4),
             price, size=12, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT, font_name=FONT_NAME)

add_footer(slide)

# ============================================================
# SLIDE 4: Professional детально
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Тариф Professional", "Расширенная аналитика для требовательных клиентов")

# Левая часть
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(0.08), Inches(5.3), ACCENT_NAVY)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.4),
         "Всё из Starter, плюс:",
         size=16, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

features_pro = [
    "Deep Diagnostic Analysis (DDA)",
    "Детекция аномалий (алгоритм Isolation Forest)",
    "Анализ сезонности через быстрое преобразование Фурье",
    "A/B анализ: сравнение периодов, оценка изменений",
    "LLM-интерпретация результатов (YandexGPT)",
    "Корреляционный анализ между тегами",
    "Прогнозирование трендов",
    "До 5 пользователей с различными ролями",
    "API доступ для интеграции с внешними системами",
    "Формирование отчётов в PDF",
]
for i, feat in enumerate(features_pro):
    add_text(slide, Inches(1.0), Inches(2.4 + i * 0.42), Inches(5.5), Inches(0.4),
             f"—  {feat}", size=11, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Правая часть — цены
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.833), Inches(2.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(0.08), Inches(2.3), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(1.8), Inches(5), Inches(0.4), "Ценообразование",
         size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

add_text(slide, Inches(7.4), Inches(2.3), Inches(5), Inches(1.4),
         "15 000 тегов:  300 000 ₽/мес  |  3 060 000 ₽/год\n"
         "50 000 тегов:  750 000 ₽/мес  |  7 650 000 ₽/год\n\n"
         "Маржа: 22-28%   |   Окупаемость для клиента: 4-6 мес",
         size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Правая часть — целевые клиенты
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(5.833), Inches(2.8), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(0.08), Inches(2.8), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(4.3), Inches(5), Inches(0.4), "Целевые сегменты",
         size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

customers_pro = [
    ("Центры обработки данных", "300 000 — 500 000 ₽/мес"),
    ("Средние производства", "300 000 — 500 000 ₽/мес"),
    ("Сети торговых центров", "500 000 — 750 000 ₽/мес"),
    ("Гостиничные комплексы", "200 000 — 350 000 ₽/мес"),
]
for i, (name, price) in enumerate(customers_pro):
    add_text(slide, Inches(7.4), Inches(4.9 + i * 0.45), Inches(3.2), Inches(0.4),
             name, size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)
    add_text(slide, Inches(10.3), Inches(4.9 + i * 0.45), Inches(2.2), Inches(0.4),
             price, size=12, color=ACCENT_NAVY, align=PP_ALIGN.RIGHT, font_name=FONT_NAME)

add_footer(slide)

# ============================================================
# SLIDE 5: Enterprise + реалистичные ожидания
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Тариф Enterprise", "Долгосрочная цель на 2-3 года, не на первый год")

# Левая часть — что входит
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(0.08), Inches(5.3), DEEP_PURPLE)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.4),
         "Всё из Professional, плюс:",
         size=16, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

features_ent = [
    "On-premise развёртывание LLM (без зависимости от внешних API)",
    "Разработка custom ML моделей под специфику клиента",
    "Мульти-объектная архитектура (неограниченно)",
    "Неограниченное количество пользователей",
    "Priority support с SLA 4 часа",
    "Выделенный Customer Success Manager",
    "Индивидуальные интеграции с ERP и MES",
    "White-label решение (собственный брендинг)",
    "Ежеквартальные бизнес-обзоры",
    "Влияние на дорожную карту продукта",
]
for i, feat in enumerate(features_ent):
    add_text(slide, Inches(1.0), Inches(2.4 + i * 0.42), Inches(5.5), Inches(0.4),
             f"—  {feat}", size=11, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Правая часть — цены и реалистичность
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.833), Inches(2.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(0.08), Inches(2.3), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(1.8), Inches(5), Inches(0.4), "Ценообразование",
         size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

add_text(slide, Inches(7.4), Inches(2.3), Inches(5), Inches(1.4),
         "150 000 тегов:  1 800 000 ₽/мес  |  18 360 000 ₽/год\n"
         "150 000+ тегов:  по запросу (от 3 000 000 ₽/мес)\n\n"
         "Маржа: 70-85% (on-premise)   |   Окупаемость: 2-3 мес",
         size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)

# Реалистичные ожидания
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(5.833), Inches(2.8), RGBColor(254, 243, 199), line_color=ACCENT_AMBER)
add_shape_bg(slide, Inches(7.0), Inches(4.1), Inches(0.08), Inches(2.8), ACCENT_AMBER)
add_text(slide, Inches(7.4), Inches(4.3), Inches(5.2), Inches(0.4),
         "Реалистичные ожидания",
         size=14, bold=True, color=ACCENT_AMBER, font_name=FONT_NAME_BOLD)

add_text(slide, Inches(7.4), Inches(4.8), Inches(5.2), Inches(2),
         "Текущий портфель SCADA: 0 клиентов Enterprise\n"
         "уровня. Первый Enterprise клиент — это цель\n"
         "на 2-3 года работы, не на первый.\n\n"
         "Целевые сегменты (долгосрочно):\n"
         "—  Крупные промышленные предприятия\n"
         "—  Объекты энергетики\n"
         "—  Сети распределённых объектов",
         size=11, color=TEXT_PRIMARY, font_name=FONT_NAME)

add_footer(slide)

# ============================================================
# SLIDE 6: Add-ons
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Дополнительные модули", "Гибкая конфигурация под конкретные потребности клиента")

addons = [
    ("Модуль энергетики",
     "Учёт потребления электроэнергии,\nработа с тарифами, прогнозы\nпотребления на основе истории",
     "+30 000 ₽/мес", ACCENT_TEAL),
    ("LLM-интерпретация",
     "Формирование нарративов на основе\nYandexGPT с включёнными расходами\nна API",
     "+40 000 ₽/мес", ACCENT_NAVY),
    ("On-premise LLM",
     "Локальная языковая модель,\nнезависимость от внешних сервисов.\nТребуется GPU-сервер",
     "+50 000 ₽/мес\n+500 000 ₽ разово", DEEP_PURPLE),
    ("Мульти-объекты",
     "Управление до 5 объектов\nв едином интерфейсе с\nконсолидированной отчётностью",
     "+50% к базовому тарифу", ACCENT_TEAL),
    ("White-label",
     "Собственный брендинг интерфейса,\nкорпоративные цвета, логотип,\nсобственный домен",
     "+20% к базовому тарифу\n+100 000 ₽ разово", ACCENT_AMBER),
    ("Priority Support",
     "Гарантированное время реакции\n4 часа, выделенный инженер,\nпрямая линия поддержки",
     "+100 000 ₽/мес", ACCENT_RED),
]

for i, (name, desc, price, color) in enumerate(addons):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.6 + row * 2.85)

    add_shape_bg(slide, left, top, Inches(3.9), Inches(2.6), BG_CARD, line_color=DIVIDER)
    add_shape_bg(slide, left, top, Inches(0.06), Inches(2.6), color)

    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.3), Inches(0.35),
             name, size=14, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

    add_shape_bg(slide, left + Inches(0.3), top + Inches(0.65), Inches(3.3), Inches(0.005), DIVIDER)

    add_text(slide, left + Inches(0.3), top + Inches(0.8), Inches(3.3), Inches(1.0), desc,
             size=11, color=TEXT_SECONDARY, font_name=FONT_NAME)

    add_shape_bg(slide, left + Inches(0.3), top + Inches(1.95), Inches(3.3), Inches(0.5), BG_DARK)
    add_text(slide, left + Inches(0.3), top + Inches(2.0), Inches(3.3), Inches(0.4), price,
             size=12, bold=False, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

add_footer(slide)

# ============================================================
# SLIDE 7: Реалистичная финмодель Год 1
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Финансовая модель: Год 1 (консервативный сценарий)",
              "Основной фокус на Starter и Professional сегменты")

# Портфель клиентов
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(0.08), Inches(5.3), ACCENT_NAVY)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.4),
         "Реалистичный клиентский портфель",
         size=16, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

clients = [
    ("Starter", "12 клиентов", "× 60 000 ₽", "720 000 ₽/мес", ACCENT_TEAL),
    ("Professional", "4 клиента", "× 350 000 ₽", "1 400 000 ₽/мес", ACCENT_NAVY),
    ("Enterprise", "0 клиентов", "—", "—", DEEP_PURPLE),
]

# Заголовки колонок
add_text(slide, Inches(1.0), Inches(2.5), Inches(1.8), Inches(0.3), "Тариф",
         size=11, bold=True, color=TEXT_SECONDARY, font_name=FONT_NAME_BOLD)
add_text(slide, Inches(2.8), Inches(2.5), Inches(1.5), Inches(0.3), "Количество",
         size=11, bold=True, color=TEXT_SECONDARY, font_name=FONT_NAME_BOLD)
add_text(slide, Inches(4.2), Inches(2.5), Inches(1), Inches(0.3), "Цена",
         size=11, bold=True, color=TEXT_SECONDARY, font_name=FONT_NAME_BOLD)
add_text(slide, Inches(5.2), Inches(2.5), Inches(1.3), Inches(0.3), "MRR",
         size=11, bold=True, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, font_name=FONT_NAME_BOLD)
add_shape_bg(slide, Inches(0.9), Inches(2.85), Inches(5.6), Inches(0.003), DIVIDER)

for i, (tier, count, price, total, color) in enumerate(clients):
    top = Inches(3.0 + i * 0.75)
    add_text(slide, Inches(1.0), top, Inches(1.8), Inches(0.4), tier,
             size=13, bold=True, color=color, font_name=FONT_NAME_BOLD)
    add_text(slide, Inches(2.8), top, Inches(1.5), Inches(0.4), count,
             size=12, color=TEXT_PRIMARY, font_name=FONT_NAME)
    add_text(slide, Inches(4.2), top, Inches(1), Inches(0.4), price,
             size=11, color=TEXT_SECONDARY, font_name=FONT_NAME)
    add_text(slide, Inches(5.2), top, Inches(1.3), Inches(0.4), total,
             size=12, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.RIGHT, font_name=FONT_NAME_BOLD)

# Итого
add_shape_bg(slide, Inches(0.9), Inches(5.3), Inches(5.6), Inches(0.003), DIVIDER)
add_shape_bg(slide, Inches(0.9), Inches(5.45), Inches(5.6), Inches(1.1), BG_DARK)
add_text(slide, Inches(1.1), Inches(5.6), Inches(2), Inches(0.35),
         "ИТОГО MRR:", size=14, bold=False, color=TEXT_LIGHT, font_name=FONT_NAME)
add_text(slide, Inches(3.5), Inches(5.6), Inches(2.8), Inches(0.35),
         "2 120 000 ₽/мес", size=16, bold=True, color=TEXT_LIGHT,
         align=PP_ALIGN.RIGHT, font_name=FONT_NAME_BOLD)
add_text(slide, Inches(1.1), Inches(6.05), Inches(5.2), Inches(0.35),
         "Годовая выручка (ARR): ~25 млн ₽",
         size=13, bold=False, color=TEXT_LIGHT, font_name=FONT_NAME)

# P&L
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.833), Inches(5.3), BG_CARD, line_color=DIVIDER)
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(0.08), Inches(5.3), ACCENT_NAVY)
add_text(slide, Inches(7.4), Inches(1.8), Inches(5.2), Inches(0.4),
         "Структура P&L — Год 1 (консервативно)",
         size=16, bold=True, color=TEXT_PRIMARY, font_name=FONT_NAME_BOLD)

pl_items = [
    ("Выручка (ARR)", "+25.4 млн ₽", True, True),
    ("", "", None, None),
    ("Расходы:", "", None, True),
    ("   Внешние API (Yandex)", "−8.5 млн ₽", False, False),
    ("   Инфраструктура", "−2.0 млн ₽", False, False),
    ("   Команда поддержки", "−4.0 млн ₽", False, False),
    ("   Продажи и маркетинг", "−6.0 млн ₽", False, False),
    ("   Разработка", "−8.0 млн ₽", False, False),
    ("   Прочие операционные", "−2.0 млн ₽", False, False),
    ("", "", None, None),
    ("Итого расходов", "−30.5 млн ₽", False, True),
]

current_top = 2.5
for item, amount, is_positive, is_bold in pl_items:
    if item == "":
        current_top += 0.15
        continue
    top = Inches(current_top)
    color_item = TEXT_PRIMARY
    color_amount = TEXT_PRIMARY

    if item == "Расходы:" or item == "Итого расходов":
        color_item = TEXT_SECONDARY
    elif item == "Выручка (ARR)":
        color_amount = ACCENT_GREEN
    elif is_positive is False:
        color_amount = ACCENT_RED

    add_text(slide, Inches(7.4), top, Inches(3.5), Inches(0.3), item,
             size=11, bold=is_bold, color=color_item, font_name=FONT_NAME)
    add_text(slide, Inches(10.9), top, Inches(1.8), Inches(0.3), amount,
             size=11, bold=is_bold, color=color_amount,
             align=PP_ALIGN.RIGHT, font_name=FONT_NAME)
    current_top += 0.38

# Итог
add_shape_bg(slide, Inches(7.2), Inches(6.1), Inches(5.4), Inches(0.7), RGBColor(254, 226, 226))
add_text(slide, Inches(7.5), Inches(6.2), Inches(2.5), Inches(0.5),
         "EBITDA Год 1 (прогноз)", size=13, bold=False, color=ACCENT_RED, font_name=FONT_NAME)
add_text(slide, Inches(10.5), Inches(6.2), Inches(1.9), Inches(0.5),
         "−5.1 млн ₽", size=16, bold=True, color=ACCENT_RED,
         align=PP_ALIGN.RIGHT, font_name=FONT_NAME_BOLD)

# Примечание
add_text(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.3),
         "Примечание: первый год — инвестиции в продукт и команду. Выход на операционную прибыль — год 2.",
         size=10, color=TEXT_SECONDARY, font_name=FONT_NAME)

add_footer(slide)

# === Сохранение ===
output_path = 'SCADA_AI_Ценовая_стратегия_ВНУТРЕННЯЯ.pptx'
prs.save(output_path)
print(f"Презентация сохранена: {output_path}")
print(f"Количество слайдов: {len(prs.slides)}")