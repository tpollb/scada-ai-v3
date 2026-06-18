from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

print('=== Генерация презентации v2: полная версия ===')
print()

try:
    # Создаём презентацию 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Цветовая схема (нейтральная, тёмная)
    BG_DARK = RGBColor(17, 24, 39)
    BG_CARD = RGBColor(31, 41, 55)
    TEXT_LIGHT = RGBColor(249, 250, 251)
    TEXT_GRAY = RGBColor(156, 163, 175)
    TEXT_DIM = RGBColor(107, 114, 128)
    ACCENT_NEUTRAL = RGBColor(107, 114, 128)
    BORDER_COLOR = RGBColor(55, 65, 81)

    CONTACTS_TEXT = "uskov-se.ru  ·  iridi.com  ·  q3mydoom@gmail.com  ·  s.uskov@iridi.tech"
    FONT_NAME = 'Inter'

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_contacts(slide):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = CONTACTS_TEXT
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DIM
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    def add_rect(slide, left, top, width, height, color, border=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if border:
            shape.line.color.rgb = BORDER_COLOR
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    def add_line(slide, x1, y1, x2, y2, color=BORDER_COLOR, width=0.5):
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
        connector.line.color.rgb = color
        connector.line.width = Pt(width)
        return connector

    def add_text(slide, text, left, top, width, height, size=18, bold=False,
                 color=TEXT_LIGHT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.alignment = align
        return txBox

    def add_bullet_list(slide, items, left, top, width, height, size=18,
                        color=TEXT_LIGHT, bullet_char="—"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{bullet_char}  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = FONT_NAME
            p.space_after = Pt(8)
        return txBox

    def add_screenshot_placeholder(slide, left, top, width, height, label):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"[ СКРИНШОТ ]\n\n{label}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        return shape

    # ========================================================================
    # СЛАЙД 1: Титульный
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.04), ACCENT_NEUTRAL)
    add_text(slide, "SCADA.AI", Inches(0.8), Inches(2.3), Inches(11.73), Inches(1.2),
             size=72, bold=True, color=TEXT_LIGHT)
    add_text(slide, "AI-ассистент для оператора промышленной SCADA-системы",
             Inches(0.8), Inches(3.5), Inches(11.73), Inches(0.8),
             size=28, bold=False, color=TEXT_GRAY)
    add_line(slide, Inches(0.8), Inches(4.5), Inches(2.8), Inches(4.5), ACCENT_NEUTRAL, 2)
    add_text(slide, "Усков Сергей Евгеньевич",
             Inches(0.8), Inches(4.9), Inches(11.73), Inches(0.6),
             size=22, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Версия 3.2.0  ·  Июнь 2026",
             Inches(0.8), Inches(5.5), Inches(11.73), Inches(0.5),
             size=14, bold=False, color=TEXT_GRAY)
    add_contacts(slide)
    print('✓ Слайд 1: Титульный')

    # ========================================================================
    # СЛАЙД 2: Суть проекта (ЗАГЛАВНЫЙ!)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Суть проекта", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.5), BG_CARD, border=True)
    add_text(slide, "Технологический прорыв на год вперёд",
             Inches(1.3), Inches(2.2), Inches(10.73), Inches(0.8),
             size=36, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_line(slide, Inches(5.5), Inches(3.2), Inches(7.83), Inches(3.2), ACCENT_NEUTRAL, 2)
    add_text(slide, "SCADA.AI — это не просто инструмент анализа.\nЭто готовый продукт, который выводит компанию\nна новый уровень автоматизации и интеллектуального управления.",
             Inches(1.3), Inches(3.5), Inches(10.73), Inches(1.5),
             size=20, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    # Три ключевых пункта
    add_rect(slide, Inches(0.8), Inches(5.7), Inches(3.7), Inches(1.0), BG_CARD, border=True)
    add_text(slide, "Готовый продукт", Inches(0.8), Inches(5.7), Inches(3.7), Inches(1.0),
             size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(4.8), Inches(5.7), Inches(3.7), Inches(1.0), BG_CARD, border=True)
    add_text(slide, "Прорыв на рынке", Inches(4.8), Inches(5.7), Inches(3.7), Inches(1.0),
             size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(8.8), Inches(5.7), Inches(3.7), Inches(1.0), BG_CARD, border=True)
    add_text(slide, "Год форы", Inches(8.8), Inches(5.7), Inches(3.7), Inches(1.0),
             size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_contacts(slide)
    print('✓ Слайд 2: Суть проекта (ЗАГЛАВНЫЙ)')

    # ========================================================================
    # СЛАЙД 3: Проблема
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Проблема", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Операторы тонут в данных",
             Inches(0.8), Inches(1.3), Inches(11.73), Inches(0.6),
             size=24, bold=False, color=TEXT_GRAY)
    items = [
        "10 000+ тегов в SCADA-системе промышленного здания",
        "Сотни аварий и событий ежедневно — физически невозможно отследить всё",
        "30-40 минут на ручной анализ одного параметра инженером",
        "Реагирование вместо предотвращения — тушим пожары, а не предупреждаем",
        "Энергоресурсы считаются «на глаз» — нет детализации затрат"
    ]
    add_bullet_list(slide, items, Inches(0.8), Inches(2.3), Inches(11.73), Inches(4.0),
                    size=22, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.5), BG_CARD, border=True)
    add_text(slide, "«Оператор не успевает думать — он только реагирует»",
             Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_contacts(slide)
    print('✓ Слайд 3: Проблема')

    # ========================================================================
    # СЛАЙД 4: Решение
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Решение", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_text(slide, "AI-ассистент, который думает за оператора",
             Inches(0.8), Inches(1.3), Inches(11.73), Inches(0.6),
             size=24, bold=False, color=TEXT_GRAY)
    items = [
        "6 независимых модулей: здоровье, энергоучёт, аналитика, логи, диалог, документация",
        "Автоматический анализ 500 000+ записей за 3-5 секунд",
        "Отчёты на естественном языке (русский) — без Excel-таблиц",
        "Интерактивные графики с прогнозами на 7/30/90/365 дней",
        "Детерминированный слой — результат всегда воспроизводим и объясним"
    ]
    add_bullet_list(slide, items, Inches(0.8), Inches(2.3), Inches(11.73), Inches(3.8),
                    size=22, color=TEXT_LIGHT)
    add_text(slide, "30 минут ручного анализа  →  5 секунд с AI",
             Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.5),
             size=28, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 4: Решение')

    # ========================================================================
    # СЛАЙД 5: Модули системы (карточки)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Модули системы", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    modules = [
        ("health", "Индекс здоровья системы", "Композитная оценка 0-100 с детализацией"),
        ("energy", "Энергоучёт", "Электричество, вода, тепло — расчёт стоимости"),
        ("analytics", "Аналитика и тренды", "Прогнозы на 7/30/90/365 дней, корреляции"),
        ("logs", "Анализ логов", "AI-анализ системных событий и ошибок"),
        ("chat", "Диалоговый интерфейс", "Вопросы на естественном языке"),
        ("docs", "Документация", "Встроенная справка и примеры использования"),
    ]
    for i, (name, title, desc) in enumerate(modules):
        row = i // 3
        col = i % 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(1.7 + row * 2.5)
        add_rect(slide, left, top, Inches(3.7), Inches(2.2), BG_CARD, border=True)
        add_text(slide, name, left + Inches(0.3), top + Inches(0.3), Inches(3.1), Inches(0.5),
                 size=20, bold=True, color=TEXT_LIGHT)
        add_line(slide, left + Inches(0.3), top + Inches(0.9), left + Inches(2.5), top + Inches(0.9),
                 BORDER_COLOR, 1)
        add_text(slide, title, left + Inches(0.3), top + Inches(1.1), Inches(3.1), Inches(0.4),
                 size=16, bold=True, color=TEXT_GRAY)
        add_text(slide, desc, left + Inches(0.3), top + Inches(1.5), Inches(3.1), Inches(0.6),
                 size=14, bold=False, color=TEXT_DIM)
    add_contacts(slide)
    print('✓ Слайд 5: Модули системы (карточки)')

    # ========================================================================
    # СЛАЙД 6: Демо — Операторский интерфейс
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Демо: операторский интерфейс", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_text(slide, "[ СКРИНШОТ ]  Home.svelte — главный экран оператора",
             Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.4),
             size=14, bold=False, color=TEXT_GRAY)
    add_screenshot_placeholder(slide,
        Inches(0.8), Inches(1.9), Inches(11.73), Inches(4.3),
        "Вставьте скриншот Home.svelte\n\nЧто должно быть видно:\n• Индекс здоровья системы (круговая диаграмма с цифрой 0-100)\n• Индекс жизнеобеспечения\n• Журнал аварий с приоритетами\n• Виджет энергозатрат\n• Чат-интерфейс справа")
    add_text(slide, "Всё что нужно оператору — на одном экране",
             Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.4),
             size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 6: Демо — операторский интерфейс')

    # ========================================================================
    # СЛАЙД 7: Демо — Интерактивная аналитика
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Демо: интерактивная аналитика", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_text(slide, "[ СКРИНШОТ ]  AnalyticsPanel — вкладка «Тренды»",
             Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.4),
             size=14, bold=False, color=TEXT_GRAY)
    add_screenshot_placeholder(slide,
        Inches(0.8), Inches(1.9), Inches(11.73), Inches(4.3),
        "Вставьте скриншот AnalyticsPanel с графиками\n\nЧто должно быть видно:\n• График температуры с линией тренда\n• Прогноз на 30 дней вперёд (пунктирная линия)\n• MA-7 — 7-дневная скользящая средняя\n• Кнопки Zoom / Download PNG\n• Переключатель периодов 7/30/90/365 дней")
    add_text(slide, "Zoom колёсиком мыши  ·  Скачивание PNG одной кнопкой  ·  Прогноз до 365 дней",
             Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.4),
             size=18, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 7: Демо — аналитика')

    # ========================================================================
    # СЛАЙД 8: Дополнительно (логи, конфигуратор, документация)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Дополнительные возможности", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    features = [
        ("Анализ системных логов",
         "AI автоматически анализирует системные события и ошибки. Выявляет аномалии, группирует проблемы, формулирует рекомендации. Оператор получает готовый отчёт вместо сырых логов."),
        ("Конфигуратор системы",
         "Полное управление модулями через веб-интерфейс. Включение/выключение модулей без перезапуска. Настройка тарифов и тегов счётчиков. Изменение промптов LLM в реальном времени."),
        ("Встроенная документация",
         "Полная справка прямо в интерфейсе. Описание модулей, API endpoints, примеры запросов к чату. Архитектура системы. История изменений (CHANGELOG). Доступ через кнопку в сайдбаре.")
    ]
    for i, (title, desc) in enumerate(features):
        left = Inches(0.8 + i * 4.0)
        add_rect(slide, left, Inches(1.6), Inches(3.7), Inches(4.5), BG_CARD, border=True)
        add_text(slide, title, left + Inches(0.3), Inches(1.9), Inches(3.1), Inches(0.6),
                 size=20, bold=True, color=TEXT_LIGHT)
        add_line(slide, left + Inches(0.3), Inches(2.6), left + Inches(2.5), Inches(2.6),
                 ACCENT_NEUTRAL, 1.5)
        add_text(slide, desc, left + Inches(0.3), Inches(2.9), Inches(3.1), Inches(3.0),
                 size=15, bold=False, color=TEXT_GRAY)
    add_text(slide, "Всё управление системой — через единый веб-интерфейс",
             Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.4),
             size=16, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 8: Дополнительные возможности')

    # ========================================================================
    # СЛАЙД 9: Экономический эффект
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Экономический эффект", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(4.0), Inches(2.5), BG_CARD, border=True)
    add_text(slide, "360×", Inches(0.8), Inches(1.6), Inches(4.0), Inches(1.5),
             size=96, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "ускорение анализа", Inches(0.8), Inches(3.0), Inches(4.0), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "30 мин  →  5 сек", Inches(0.8), Inches(3.4), Inches(4.0), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Экономия расходов:", Inches(5.3), Inches(1.5), Inches(7.0), Inches(0.5),
             size=20, bold=True, color=TEXT_LIGHT)
    savings = [
        "Время инженеров-аналитиков: -30 часов/месяц на отчёты",
        "Предотвращение аварий через прогнозирование (до 365 дней)",
        "Оптимизация энергопотребления через интервальные тарифы",
        "Раннее выявление битых датчиков (экономия на ложных срабатываниях)"
    ]
    add_bullet_list(slide, savings, Inches(5.3), Inches(2.1), Inches(7.0), Inches(2.2),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_text(slide, "Повышение эффективности:", Inches(5.3), Inches(4.4), Inches(7.0), Inches(0.5),
             size=20, bold=True, color=TEXT_LIGHT)
    revenue = [
        "Оператор принимает решения на основе данных, а не интуиции",
        "Прогнозы на 90/365 дней позволяют планировать ремонты",
        "Прозрачная отчётность для руководства (PDF одним кликом)"
    ]
    add_bullet_list(slide, revenue, Inches(5.3), Inches(5.0), Inches(7.0), Inches(1.4),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_rect(slide, Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5), BG_CARD, border=True)
    add_text(slide, "Экономит расходы  +  Повышает эффективность операций  =  Максимальный тип предложения",
             Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_contacts(slide)
    print('✓ Слайд 9: Экономический эффект')

    # ========================================================================
    # СЛАЙД 10: Стадия готовности
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Стадия готовности", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(4.0), Inches(2.8), BG_CARD, border=True)
    add_text(slide, "85%", Inches(0.8), Inches(1.8), Inches(4.0), Inches(2.0),
             size=120, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "готовность к production", Inches(0.8), Inches(3.7), Inches(4.0), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "Выпуск за 1 неделю", Inches(0.8), Inches(4.2), Inches(4.0), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Что сделано:", Inches(5.3), Inches(1.8), Inches(7.0), Inches(0.5),
             size=20, bold=True, color=TEXT_LIGHT)
    done_items = [
        "Рабочая версия 3.2.0 с 6 модулями",
        "Интеграция с реальной БД SCADA (PostgreSQL / TimescaleDB)",
        "Обработка 500 000+ записей в production-режиме",
        "Интерактивная визуализация (Chart.js + zoom/pan)",
        "Русифицированный UI со всеми виджетами",
        "Полная документация (6 файлов markdown)",
        "История версий с v1.x.x и CHANGELOG"
    ]
    add_bullet_list(slide, done_items, Inches(5.3), Inches(2.4), Inches(7.0), Inches(4.0),
                    size=16, color=TEXT_LIGHT, bullet_char="✓")
    add_rect(slide, Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5), BG_CARD, border=True)
    add_text(slide, "Идея  ·  Прототип  ·  Тестирование  ·  Готов к реализации  ·  В продакшене",
             Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_contacts(slide)
    print('✓ Слайд 10: Стадия готовности')

    # ========================================================================
    # СЛАЙД 11: Планы развития
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Планы развития", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), BG_CARD, border=True)
    add_text(slide, "v3.3.0  ·  Deep Data Analysis", Inches(1.1), Inches(1.7), Inches(5.0), Inches(0.6),
             size=22, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Хирургический инструмент для инженеров-аналитиков",
             Inches(1.1), Inches(2.3), Inches(5.0), Inches(0.5),
             size=14, bold=False, color=TEXT_GRAY)
    v33_items = [
        "Isolation Forest — детекция аномалий",
        "FFT — обнаружение сезонных циклов",
        "Granger Causality — причинно-следственные связи",
        "A/B сравнение периодов (t-test, Mann-Whitney)",
        "Экспорт отчётов в PDF и Excel",
        "Выбор любого тега или группы тегов"
    ]
    add_bullet_list(slide, v33_items, Inches(1.1), Inches(2.9), Inches(5.0), Inches(3.0),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_rect(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5), BG_CARD, border=True)
    add_text(slide, "v4.0.0  ·  Долгосрочные цели", Inches(7.3), Inches(1.7), Inches(5.0), Inches(0.6),
             size=22, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Расширение аудитории и возможностей",
             Inches(7.3), Inches(2.3), Inches(5.0), Inches(0.5),
             size=14, bold=False, color=TEXT_GRAY)
    v40_items = [
        "Boss Dashboard — экран для руководства (KPI)",
        "Ролевая модель (admin / engineer / operator / boss)",
        "Mobile App — доступ с планшета",
        "Multi-tenancy — масштабирование на здания",
        "Telegram-интеграция для алертов",
        "Расписания проверок с уведомлениями"
    ]
    add_bullet_list(slide, v40_items, Inches(7.3), Inches(2.9), Inches(5.0), Inches(3.0),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_text(slide, "Цель: превратить SCADA.AI в полноценный BI-инструмент для промышленного предприятия",
             Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 11: Планы развития')

    # ========================================================================
    # СЛАЙД 12: Грант Yandex.AI.Boost + SaaS (ОБНОВЛЁН!)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Поддержка и коммерциализация", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.8), BG_CARD, border=True)
    add_text(slide, "Yandex.AI.Boost", Inches(1.3), Inches(1.8), Inches(10.73), Inches(0.8),
             size=36, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Грантовая программа поддержки AI-проектов",
             Inches(1.3), Inches(2.6), Inches(10.73), Inches(0.5),
             size=18, bold=False, color=TEXT_GRAY)
    add_line(slide, Inches(1.3), Inches(3.3), Inches(12.03), Inches(3.3), BORDER_COLOR, 1)
    add_text(slide, "Текущий этап:", Inches(1.3), Inches(3.6), Inches(5.0), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT)
    current_items = [
        "Грант 100 000 ₽ получен и освоен",
        "Используется YandexGPT 5.1 как основная LLM",
        "Еженедельная отчётность перед Яндексом",
        "Все API-запросы идут через Яндекс.Облако"
    ]
    add_bullet_list(slide, current_items, Inches(1.3), Inches(4.2), Inches(5.0), Inches(2.0),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_text(slide, "Следующий этап:", Inches(7.0), Inches(3.6), Inches(5.0), Inches(0.5),
             size=17, bold=True, color=TEXT_LIGHT)
    next_items = [
        "Увеличение гранта до 1 000 000 ₽",
        "Размещение в Yandex Cloud Marketplace",
        "SaaS-модель продаж (подписка)",
        "Продакшн-внедрение на первом объекте"
    ]
    add_bullet_list(slide, next_items, Inches(7.0), Inches(4.2), Inches(5.0), Inches(2.0),
                    size=15, color=TEXT_LIGHT, bullet_char="—")
    add_rect(slide, Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5), BG_CARD, border=True)
    add_text(slide, "Проведены консультации с Яндекс по модели продаж SaaS. Достигнуты договорённости о размещении на маркетплейсе приложений.",
             Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.5),
             size=15, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_contacts(slide)
    print('✓ Слайд 12: Yandex.AI.Boost + SaaS')

    # ========================================================================
    # СЛАЙД 13: Ключевые преимущества
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Ключевые преимущества", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    advantages = [
        ("Детерминированность",
         "Формулы вместо «чёрного ящика». Каждый штраф, каждый индекс — объясним. Оператор видит математику, а не магию. Полная воспроизводимость результатов."),
        ("Скорость",
         "500 000 записей анализируются за 3-5 секунд. Python-агрегация + LLM-интерпретация. Быстрее любого ручного Excel-отчёта. Мгновенная обратная связь оператору."),
        ("Расширяемость",
         "Новый модуль = 1 папка + 3 файла. Добавить анализ нового параметра — 1 день работы. Модульная архитектура не требует переписывания системы при расширении.")
    ]
    for i, (title, desc) in enumerate(advantages):
        left = Inches(0.8 + i * 4.0)
        add_rect(slide, left, Inches(1.6), Inches(3.8), Inches(4.2), BG_CARD, border=True)
        add_text(slide, title, left + Inches(0.3), Inches(1.9), Inches(3.2), Inches(0.6),
                 size=22, bold=True, color=TEXT_LIGHT)
        add_line(slide, left + Inches(0.3), Inches(2.6), left + Inches(2.5), Inches(2.6),
                 ACCENT_NEUTRAL, 1.5)
        add_text(slide, desc, left + Inches(0.3), Inches(2.9), Inches(3.2), Inches(2.7),
                 size=15, bold=False, color=TEXT_GRAY)
    add_text(slide, "Интерактивная визуализация  ·  Zoom/pan колёсиком мыши  ·  Экспорт PNG / PDF / Excel",
             Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.4),
             size=16, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 13: Ключевые преимущества')

    # ========================================================================
    # СЛАЙД 14: Технологический стек (обобщённо, в конце)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.04), Inches(7.5), ACCENT_NEUTRAL)
    add_text(slide, "Технологический стек", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
             size=40, bold=True, color=TEXT_LIGHT)
    add_text(slide, "Современные, проверенные решения",
             Inches(0.8), Inches(1.3), Inches(11.73), Inches(0.6),
             size=20, bold=False, color=TEXT_GRAY)
    add_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0), BG_CARD, border=True)
    add_text(slide, "Backend", Inches(1.3), Inches(2.5), Inches(4.5), Inches(0.5),
             size=20, bold=True, color=TEXT_LIGHT)
    backend_items = [
        "Python 3.13 — современный язык",
        "Асинхронный HTTP-сервер",
        "Реляционная БД с расширениями для time-series",
        "YandexGPT 5.1 — большая языковая модель",
        "NumPy / SciPy — статистическая обработка",
        "Chart.js — серверная генерация графиков"
    ]
    add_bullet_list(slide, backend_items, Inches(1.3), Inches(3.1), Inches(4.5), Inches(3.0),
                    size=15, color=TEXT_LIGHT, bullet_char="·")
    add_rect(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0), BG_CARD, border=True)
    add_text(slide, "Frontend", Inches(7.5), Inches(2.5), Inches(4.5), Inches(0.5),
             size=20, bold=True, color=TEXT_LIGHT)
    frontend_items = [
        "Современный реактивный фреймворк (Svelte 5)",
        "Utility-first стилизация (Tailwind CSS)",
        "Интерактивные графики с zoom/pan",
        "Адаптивный дизайн (mobile + desktop)",
        "Минимальный размер бандла",
        "Быстрая hot-reload разработка"
    ]
    add_bullet_list(slide, frontend_items, Inches(7.5), Inches(3.1), Inches(4.5), Inches(3.0),
                    size=15, color=TEXT_LIGHT, bullet_char="·")
    add_text(slide, "Модульная архитектура: добавление нового модуля = 1 папка + 3 файла",
             Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.4),
             size=16, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_contacts(slide)
    print('✓ Слайд 14: Технологический стек')

    # ========================================================================
    # СЛАЙД 15: Спасибо / Контакты (финальный)
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.04), ACCENT_NEUTRAL)
    add_text(slide, "Спасибо за внимание",
             Inches(0.8), Inches(2.2), Inches(11.73), Inches(1.0),
             size=56, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "Готов ответить на вопросы и показать живое демо",
             Inches(0.8), Inches(3.4), Inches(11.73), Inches(0.6),
             size=24, bold=False, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_line(slide, Inches(5.5), Inches(4.3), Inches(7.83), Inches(4.3), ACCENT_NEUTRAL, 2)
    add_text(slide, "Усков Сергей Евгеньевич",
             Inches(0.8), Inches(4.7), Inches(11.73), Inches(0.6),
             size=22, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "uskov-se.ru", Inches(2.5), Inches(5.4), Inches(3.5), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, "iridi.com", Inches(7.5), Inches(5.4), Inches(3.5), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "q3mydoom@gmail.com", Inches(2.5), Inches(5.8), Inches(3.5), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, "s.uskov@iridi.tech", Inches(7.5), Inches(5.8), Inches(3.5), Inches(0.4),
             size=15, bold=False, color=TEXT_GRAY, align=PP_ALIGN.LEFT)
    add_contacts(slide)
    print('✓ Слайд 15: Спасибо + контакты')

    # ========================================================================
    # СОХРАНЕНИЕ
    # ========================================================================
    output_path = Path('Усков_SCADA_AI_Dark.pptx')
    prs.save(output_path)

    print()
    print('=' * 70)
    print(f'✓ ПРЕЗЕНТАЦИЯ СОХРАНЕНА: {output_path.absolute()}')
    print('=' * 70)
    print()
    print(f'Файл находится в: {Path.cwd()}')
    print(f'Размер: {output_path.stat().st_size / 1024:.1f} KB')
    print()
    print('СТРУКТУРА (15 слайдов):')
    print('  1. Титульный')
    print('  2. Суть проекта (ЗАГЛАВНЫЙ — прорыв на год)')
    print('  3. Проблема')
    print('  4. Решение')
    print('  5. Модули (карточки)')
    print('  6. Демо: операторский интерфейс [СКРИНШОТ]')
    print('  7. Демо: аналитика [СКРИНШОТ]')
    print('  8. Дополнительно: логи, конфигуратор, документация')
    print('  9. Экономический эффект (360×)')
    print(' 10. Стадия готовности (85%, 1 неделя до прода)')
    print(' 11. Планы развития (v3.3.0 + v4.0.0)')
    print(' 12. Yandex.AI.Boost + SaaS маркетплейс')
    print(' 13. Ключевые преимущества')
    print(' 14. Технологический стек (обобщённо)')
    print(' 15. Спасибо + контакты')

except Exception as e:
    print()
    print('❌ ОШИБКА при генерации:')
    print(f'   {type(e).__name__}: {e}')
    import traceback
    traceback.print_exc()