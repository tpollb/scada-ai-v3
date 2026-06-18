from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

print('=== Генерация презентации Усков_SCADA_AI.pptx ===')
print()

# Создаём презентацию 16:9
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Цвета SCADA.AI
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(31, 41, 55)      # #1F2937 — основной текст
MEDIUM_GRAY = RGBColor(107, 114, 128)  # #6B7280 — подзаголовки
LIGHT_GRAY = RGBColor(204, 204, 204)   # #CCCCCC — контакты
ACCENT_BLUE = RGBColor(37, 99, 235)    # #2563EB — акцент
ACCENT_GREEN = RGBColor(22, 163, 74)   # #16A34A
ACCENT_RED = RGBColor(220, 38, 38)     # #DC2626
BG_LIGHT = RGBColor(249, 250, 251)     # #F9FAFB

# Контакты (появятся внизу каждого слайда)
CONTACTS_TEXT = "uskov-se.ru  |  iridi.com  |  q3mydoom@gmail.com  |  s.uskov@iridi.tech"

def add_contacts(slide):
    """Добавить контакты внизу слайда (еле различимые)"""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12.33)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = CONTACTS_TEXT
    p.font.size = Pt(8)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

def add_title_text(slide, text, left, top, width, height, size=32, bold=True, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Добавить текстовый блок"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    p.alignment = align
    return txBox

def add_body_text(slide, text, left, top, width, height, size=18, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Добавить основной текст"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    p.alignment = align
    return txBox

def add_bullet_list(slide, items, left, top, width, height, size=18, color=DARK_GRAY, bullet_color=ACCENT_BLUE):
    """Добавить маркированный список"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        p.level = 0
        # Используем символ • как маркер (надёжнее чем bullet)
        p.text = f"●  {item}"
    return txBox

def add_rect(slide, left, top, width, height, color):
    """Добавить прямоугольник (акцентная полоска)"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_screenshot_placeholder(slide, left, top, width, height, label):
    """Добавить плейсхолдер для скриншота"""
    shape = slide.shapes.add_shape(
        1,  # Прямоугольник
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_LIGHT
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(2)
    
    # Текст внутри плейсхолдера
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"📷 СКРИНШОТ\n\n{label}"
    p.font.size = Pt(16)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ============================================================================
# СЛАЙД 1: Титульный
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой layout
# Синяя полоска сверху
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), ACCENT_BLUE)

# Заголовок
add_title_text(slide, "SCADA.AI", Inches(0.8), Inches(2.0), Inches(11.73), Inches(1.2),
               size=72, bold=True, color=DARK_GRAY)
# Подзаголовок
add_title_text(slide, "AI-ассистент для оператора промышленной SCADA-системы",
               Inches(0.8), Inches(3.2), Inches(11.73), Inches(0.8),
               size=28, bold=False, color=MEDIUM_GRAY)

# Линия-разделитель
add_rect(slide, Inches(0.8), Inches(4.2), Inches(2.0), Inches(0.04), ACCENT_BLUE)

# Автор
add_title_text(slide, "Усков Сергей Евгеньевич",
               Inches(0.8), Inches(4.6), Inches(11.73), Inches(0.6),
               size=22, bold=True, color=DARK_GRAY)

# Версия
add_title_text(slide, "Версия 3.2.0  •  Июнь 2026",
               Inches(0.8), Inches(5.2), Inches(11.73), Inches(0.5),
               size=16, bold=False, color=MEDIUM_GRAY)

add_contacts(slide)


# ============================================================================
# СЛАЙД 2: Проблема
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Проблема", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)
add_title_text(slide, "Операторы тонут в данных", Inches(0.8), Inches(1.3), Inches(11.73), Inches(0.6),
               size=24, bold=False, color=MEDIUM_GRAY)

items = [
    "10 000+ тегов в SCADA-системе промышленного здания",
    "Сотни аварий и событий ежедневно — физически невозможно отследить всё",
    "30-40 минут на ручной анализ одного параметра инженером",
    "Реагирование вместо предотвращения — тушим пожары, а не предупреждаем",
    "Энергоресурсы считаются «на глаз» — нет детализации затрат"
]
add_bullet_list(slide, items, Inches(0.8), Inches(2.3), Inches(11.73), Inches(4.5), size=22)

# Цитата-акцент внизу
add_title_text(slide, "«Оператор не успевает думать — он только реагирует»",
               Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.6),
               size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 3: Решение
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Решение", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)
add_title_text(slide, "AI-ассистент, который думает за оператора",
               Inches(0.8), Inches(1.3), Inches(11.73), Inches(0.6),
               size=24, bold=False, color=MEDIUM_GRAY)

items = [
    "5 независимых модулей — здоровье, энергоучёт, аналитика, логи, диалог",
    "Автоматический анализ 500 000+ записей за 3-5 секунд",
    "Отчёты на естественном языке (русский) — не нужны Excel-таблицы",
    "Интерактивные графики с прогнозами на 7/30/90/365 дней",
    "Детерминированный слой — результат всегда воспроизводим и объясним"
]
add_bullet_list(slide, items, Inches(0.8), Inches(2.3), Inches(11.73), Inches(4.0), size=22)

# Финальная фраза-удар
add_title_text(slide, "30 минут ручного анализа  →  5 секунд с AI",
               Inches(0.8), Inches(6.0), Inches(11.73), Inches(0.6),
               size=28, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 4: Модули системы
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Модули системы", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# Создаём таблицу модулей
modules = [
    ("health", "Индекс здоровья системы (0-100)", "HealthScoreCard, LifeSupportCard"),
    ("energy_electricity", "Расчёт стоимости электроэнергии", "EnergyCostCard"),
    ("energy_water", "Учёт потребления воды", "EnergyCostCard"),
    ("energy_heat", "Учёт потребления тепла", "EnergyCostCard"),
    ("analytics", "Тренды, прогнозы, корреляции", "AnalyticsPanel (4 вкладки)"),
    ("logs", "AI-анализ системных логов", "NarrativePanel"),
]

# Создаём таблицу
rows = len(modules) + 1
cols = 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.5))
table = table_shape.table

# Ширина колонок
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(5.5)
table.columns[2].width = Inches(3.43)

# Заголовки
headers = ["Модуль", "Назначение", "Виджеты"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

# Данные
for i, (name, desc, widgets) in enumerate(modules, 1):
    row_data = [name, desc, widgets]
    for j, val in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Segoe UI'
        if j == 0:
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
        # Чередование фона
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT

add_contacts(slide)


# ============================================================================
# СЛАЙД 5: Демо - Операторский интерфейс
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Демо: операторский интерфейс", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

add_title_text(slide, "📷  СКРИНШОТ: Home.svelte (главный экран оператора)",
               Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.5),
               size=16, bold=False, color=MEDIUM_GRAY)

# Плейсхолдер для скриншота
add_screenshot_placeholder(slide,
    Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5),
    "Вставьте скриншот Home.svelte\n\nВидно:\n• Индекс здоровья системы (круг с цифрой)\n• Индекс жизнеобеспечения\n• Журнал аварий\n• Виджет энергозатрат\n• Чат-интерфейс справа")

add_title_text(slide, "Всё что нужно оператору — на одном экране",
               Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.5),
               size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 6: Демо - Аналитика
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Демо: интерактивная аналитика", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

add_title_text(slide, "📷  СКРИНШОТ: AnalyticsPanel (вкладка «Тренды»)",
               Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.5),
               size=16, bold=False, color=MEDIUM_GRAY)

add_screenshot_placeholder(slide,
    Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5),
    "Вставьте скриншот AnalyticsPanel с графиками\n\nВидно:\n• График температуры с линией тренда\n• Прогноз на 30 дней вперёд\n• MA-7 (скользящая средняя)\n• Кнопки Zoom / Download PNG\n• Переключатель периодов 7/30/90/365 дней")

add_title_text(slide, "Zoom колёсиком мыши  •  Скачивание PNG одной кнопкой  •  Прогноз до 365 дней",
               Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.5),
               size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 7: Технологический стек
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Технологический стек", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# Backend
add_title_text(slide, "Backend", Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
               size=26, bold=True, color=ACCENT_BLUE)
backend_items = [
    "FastAPI — высокопроизводительный REST API",
    "asyncpg + PostgreSQL (TimescaleDB) — хранилище метрик",
    "YandexGPT 5.1 — LLM с tool calling",
    "NumPy / SciPy — статистическая обработка",
    "Chart.js — серверная генерация графиков для PDF"
]
add_bullet_list(slide, backend_items, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0), size=16)

# Frontend
add_title_text(slide, "Frontend", Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6),
               size=26, bold=True, color=ACCENT_BLUE)
frontend_items = [
    "Svelte 5 с runes — современный реактивный UI",
    "Tailwind CSS — utility-first стилизация",
    "Chart.js + svelte-chartjs — интерактивные графики",
    "chartjs-plugin-zoom — масштабирование колёсиком",
    "ky — лёгкий HTTP-клиент"
]
add_bullet_list(slide, frontend_items, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0), size=16)

# Итоговая строка
add_title_text(slide, "Модульная архитектура: добавление нового модуля = 1 папка + 3 файла",
               Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.6),
               size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 8: Экономический эффект
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Экономический эффект", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# Большая цифра экономии времени
add_title_text(slide, "360×", Inches(0.8), Inches(1.5), Inches(4.0), Inches(1.5),
               size=96, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_title_text(slide, "ускорение анализа", Inches(0.8), Inches(2.9), Inches(4.0), Inches(0.5),
               size=18, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_title_text(slide, "30 мин → 5 сек", Inches(0.8), Inches(3.3), Inches(4.0), Inches(0.4),
               size=16, bold=False, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Правый блок с пунктами
add_title_text(slide, "Экономия расходов:", Inches(5.3), Inches(1.5), Inches(7.0), Inches(0.5),
               size=20, bold=True, color=DARK_GRAY)
savings = [
    "Время инженеров-аналитиков: -30 часов/месяц на отчёты",
    "Предотвращение аварий через прогнозирование (прогноз до 365 дней)",
    "Оптимизация энергопотребления через интервальные тарифы",
    "Раннее выявление битых датчиков (экономия на ложных срабатываниях)"
]
add_bullet_list(slide, savings, Inches(5.3), Inches(2.1), Inches(7.0), Inches(2.2), size=15)

add_title_text(slide, "Повышение эффективности:", Inches(5.3), Inches(4.5), Inches(7.0), Inches(0.5),
               size=20, bold=True, color=DARK_GRAY)
revenue = [
    "Оператор принимает решения на основе данных, а не интуиции",
    "Прогнозы на 90/365 дней позволяют планировать ремонты",
    "Прозрачная отчётность для руководства (PDF одним кликом)"
]
add_bullet_list(slide, revenue, Inches(5.3), Inches(5.1), Inches(7.0), Inches(1.6), size=15)

# Итоговая оценка
add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.6), ACCENT_BLUE)
add_title_text(slide, "Тип 3: Экономит расходы + Повышает доходы  =  3 балла",
               Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.6),
               size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 9: Стадия готовности
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Стадия готовности", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# Большая цифра
add_title_text(slide, "4", Inches(0.8), Inches(1.8), Inches(4.0), Inches(2.0),
               size=144, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_title_text(slide, "балла", Inches(0.8), Inches(3.6), Inches(4.0), Inches(0.6),
               size=28, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_title_text(slide, "Прототип готов\nк реализации", Inches(0.8), Inches(4.2), Inches(4.0), Inches(1.0),
               size=18, bold=False, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Чеклист справа
add_title_text(slide, "Что сделано:", Inches(5.3), Inches(1.8), Inches(7.0), Inches(0.5),
               size=22, bold=True, color=DARK_GRAY)
done_items = [
    "✓  Рабочая версия 3.2.0 с 6 модулями",
    "✓  Интеграция с реальной БД SCADA (PostgreSQL / TimescaleDB)",
    "✓  Обработка 500 000+ записей в production-режиме",
    "✓  Интерактивная визуализация (Chart.js + zoom/pan)",
    "✓  Русифицированный UI со всеми виджетами",
    "✓  Полная документация (6 файлов markdown)",
    "✓  История версий и CHANGELOG с v1.x.x"
]
add_bullet_list(slide, done_items, Inches(5.3), Inches(2.4), Inches(7.0), Inches(4.0), size=16, color=DARK_GRAY)

# Шкала внизу
add_title_text(slide, "Идея  →  Прототип  →  Тестирование  →  Готов к реализации  →  Запущен",
               Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.5),
               size=16, bold=False, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_title_text(slide, "      1                    2                    3                       4 ★                  5",
               Inches(0.8), Inches(6.7), Inches(11.73), Inches(0.4),
               size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 10: Roadmap и планы развития
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Планы развития", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# v3.3.0 — ближайший релиз
add_title_text(slide, "v3.3.0 — Deep Data Analysis (в разработке)",
               Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
               size=22, bold=True, color=ACCENT_BLUE)
v33_items = [
    "Isolation Forest — детекция аномалий",
    "FFT — обнаружение сезонных циклов (сутки/неделя)",
    "Granger Causality — причинно-следственные связи",
    "A/B сравнение периодов (t-test, Mann-Whitney)",
    "Экспорт отчётов в PDF и Excel"
]
add_bullet_list(slide, v33_items, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.0), size=15)

# v4.0.0 — долгосрочно
add_title_text(slide, "v4.0.0 — долгосрочные цели",
               Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.6),
               size=22, bold=True, color=ACCENT_BLUE)
v40_items = [
    "Boss Dashboard — отдельный экран для руководства (KPI, метрики)",
    "Ролевая модель доступа (admin / engineer / operator / boss)",
    "Mobile App (React Native) — доступ с планшета",
    "Multi-tenancy — масштабирование на несколько зданий",
    "Интеграция с Telegram для алертов"
]
add_bullet_list(slide, v40_items, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.0), size=15)

# Итоговая строка
add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.73), Inches(0.8), BG_LIGHT)
add_title_text(slide, "Цель: превратить SCADA.AI в полноценный BI-инструмент для промышленного предприятия",
               Inches(0.8), Inches(5.7), Inches(11.73), Inches(0.5),
               size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 11: Ключевые преимущества
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_title_text(slide, "Ключевые преимущества", Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8),
               size=40, bold=True, color=DARK_GRAY)

# Три колонки с преимуществами
col_data = [
    ("🎯", "Детерминированность", "Формулы вместо «чёрного ящика». Каждый штраф, каждый индекс — объясним. Оператор видит математику, а не магию."),
    ("⚡", "Скорость", "500 000 записей анализируются за 3-5 секунд. Python-агрегация + LLM-интерпретация. Быстрее любого Excel-отчёта."),
    ("🔌", "Расширяемость", "Новый модуль = 1 папка + 3 файла. Добавить анализ нового параметра — 1 день работы. Не надо переписывать систему.")
]

for i, (icon, title, desc) in enumerate(col_data):
    left = Inches(0.8 + i * 4.0)
    # Иконка
    add_title_text(slide, icon, left, Inches(1.6), Inches(3.8), Inches(0.8),
                   size=48, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    # Заголовок
    add_title_text(slide, title, left, Inches(2.4), Inches(3.8), Inches(0.6),
                   size=22, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    # Описание
    add_body_text(slide, desc, left, Inches(3.1), Inches(3.8), Inches(3.0),
                  size=15, color=MEDIUM_GRAY, align=PP_ALIGN.LEFT)

# Ещё одна строка преимуществ
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.73), Inches(0.8), BG_LIGHT)
add_title_text(slide, "Интерактивная визуализация  •  Zoom/pan колёсиком мыши  •  Экспорт PNG/PDF/Excel",
               Inches(0.8), Inches(6.0), Inches(11.73), Inches(0.5),
               size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_contacts(slide)


# ============================================================================
# СЛАЙД 12: Спасибо / Контакты (финальный)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), ACCENT_BLUE)

add_title_text(slide, "Спасибо за внимание",
               Inches(0.8), Inches(2.0), Inches(11.73), Inches(1.0),
               size=56, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_title_text(slide, "Готов ответить на вопросы и показать живое демо",
               Inches(0.8), Inches(3.2), Inches(11.73), Inches(0.6),
               size=24, bold=False, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Разделитель
add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.33), Inches(0.04), ACCENT_BLUE)

# Контакты крупно
add_title_text(slide, "Усков Сергей Евгеньевич",
               Inches(0.8), Inches(4.6), Inches(11.73), Inches(0.6),
               size=22, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

contacts_large = [
    "🌐  uskov-se.ru",
    "🏢  iridi.com",
    "✉   q3mydoom@gmail.com",
    "✉   s.uskov@iridi.tech"
]
add_bullet_list(slide, contacts_large, Inches(4.0), Inches(5.3), Inches(5.33), Inches(2.0),
                size=16, color=DARK_GRAY)

add_contacts(slide)


# ============================================================================
# Сохраняем
# ============================================================================
output_path = Path('Усков_SCADA_AI.pptx')
prs.save(output_path)

print(f'✓ Презентация сохранена: {output_path.absolute()}')
print()
print('=' * 70)
print('СТРУКТУРА ПРЕЗЕНТАЦИИ (12 слайдов, ~6-7 минут питча):')
print('=' * 70)
print()
print('Слайд 1:  Титульный (SCADA.AI, Усков С.Е.)')
print('Слайд 2:  Проблема — операторы тонут в данных')
print('Слайд 3:  Решение — AI-ассистент с 5 модулями')
print('Слайд 4:  Модули системы (таблица)')
print('Слайд 5:  📷 Демо: операторский интерфейс (Home.svelte)')
print('Слайд 6:  📷 Демо: интерактивная аналитика (AnalyticsPanel)')
print('Слайд 7:  Технологический стек (Backend + Frontend)')
print('Слайд 8:  Экономический эффект (360× ускорение, тип 3 = 3 балла)')
print('Слайд 9:  Стадия готовности (4 балла — готов к реализации)')
print('Слайд 10: Планы развития (v3.3.0 Deep Analysis + v4.0.0)')
print('Слайд 11: Ключевые преимущества (3 колонки)')
print('Слайд 12: Спасибо + крупные контакты')
print()
print('=' * 70)
print('ЧТО СДЕЛАТЬ ВРУЧНУЮ:')
print('=' * 70)
print()
print('1. Открой Усков_SCADA_AI.pptx в PowerPoint')
print()
print('2. На Слайде 5 замени плейсхолдер на скриншот:')
print('   → Home.svelte (главный экран оператора)')
print('   → Должно быть видно: индекс здоровья, журналы аварий, чат')
print()
print('3. На Слайде 6 замени плейсхолдер на скриншот:')
print('   → AnalyticsPanel во вкладке "Тренды"')
print('   → Должно быть видно: график с линией тренда и прогнозом')
print()
print('4. Опционально: добавь логотип компании на Слайд 1 (правый верхний угол)')
print()
print('5. Проверь итоговый размер: ~12 слайдов × 30 сек = 6 минут питча')
print()
print('=' * 70)
print('ИТОГОВАЯ ОЦЕНКА ПО КРИТЕРИЯМ КОНКУРСА:')
print('=' * 70)
print()
print('✓ Практическая применимость: рабочая версия на реальной БД')
print('✓ Экономический эффект: тип 3 = 3 балла (экономит + повышает)')
print('✓ Стадия готовности: 4 балла (прототип готов к реализации)')
print()
print('ПРОГНОЗ: 7 баллов из возможных. Сильная заявка! 🎯')